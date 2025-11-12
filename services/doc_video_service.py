# -*- coding: utf-8 -*-
"""
Doc → Video Service (thematic, robust, fast)

- Ảnh phải bám nội dung: heading + keywords của slide.
- Chỉ dùng ảnh thật/bối cảnh (lọc icon/clipart/logo/mockup).
- Có “cứu cánh” khi lượt lọc đầu quá gắt.
- Tối ưu tốc độ: time budget/slide, timeout tải ảnh ngắn, bỏ zoom theo thời gian,
  encode fps thấp hơn và preset nhanh, TTS fallback im lặng khi không có human voice.

Cách dùng (ví dụ):
    # python -c 'from doc_video_service import make_video_from_file as run; print(run("input.pdf","out"))'
"""
from __future__ import annotations

import os, io, uuid, textwrap, re, gc, shutil, tempfile, contextlib, time, math, random
from typing import List, Tuple, Dict, Any, Optional

from io import BytesIO
from PIL import Image, ImageDraw, ImageFont, ImageOps, ImageFilter
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, CompositeVideoClip
from moviepy.video.fx import fadein, fadeout
from moviepy.audio.AudioClip import AudioArrayClip, concatenate_audioclips
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed

import requests
import numpy as np

# === External services (có sẵn trong dự án) ===
from services.human_voice_service import is_human_voice_available, synthesize_human_voice
from services.image_search_service import search_images

# === CẤU HÌNH CHUNG ===
MAX_SLIDES = 50
MAX_CHARS_PER_CHUNK = 6000
VIDEO_WIDTH, VIDEO_HEIGHT = 1280, 720
WPM = 140
MAX_TEXT_LENGTH = 500000          # 500k ký tự
MAX_DURATION_SEC = 600

# Chất lượng & đa dạng ảnh
MIN_IMG_W, MIN_IMG_H = 900, 600   # hạ nhẹ để bớt loại nhầm ảnh gốc tốt
FINAL_IMAGES_PER_SLIDE = 3        # 2-3 ảnh/slide (sẽ randomize)
MAX_FETCH_PER_QUERY = 20          # Tăng để đa dạng và tránh fallback
MAX_QUERIES_PER_SLIDE = 6
HASH_SIZE = 8
MIN_HASH_DISTANCE = 15            # Tăng để đảm bảo đa dạng rõ rệt

# Tối ưu tốc độ
ENCODE_FPS = 18                  # encode nhanh hơn, vẫn mượt
PER_IMAGE_HTTP_TIMEOUT = 4       # tải ảnh >4s bỏ qua
PER_SLIDE_TIME_BUDGET = 15.0     # Tăng: 15s/slide để tìm nhiều hơn, tránh fallback

SKIP_GTTS_IF_NO_HUMAN = True     # không có human-voice thì bỏ gTTS (tránh chờ mạng)

# Neo chủ đề (giữ tính liên quan)
TOPIC_ANCHORS = [
    "education", "learning", "student", "classroom", "elearning", "school",
    "information technology", "it", "technology", "computer", "programming",
    "coding", "software", "ai", "artificial intelligence", "machine learning",
    "data science", "robotics", "foreign language", "language learning", "english", "esl"
]

# Loại trừ chủ đề lạc hướng
BANNED_KEYWORDS = ["sports", "food", "travel", "fashion", "celebrity", "car", "weapon", "politics"]

STOPWORDS = set("""
a an the and or for with of in on to from by is are was were be been being as at that this
these those there here it its their his her your our not no do does did done using use about
""".split())

# Yêu cầu giọng người thật?
REQUIRE_HUMAN_VOICE_ONLY = os.getenv("REQUIRE_HUMAN_VOICE", "0").lower() in ("1", "true", "yes")

# Hạn chế từ khóa xấu khi truy vấn ảnh - LOẠI BỎ ẢNH CÓ NGƯỜI
NEGATIVE_PROMPT = " -icon -icons -logo -vector -svg -clipart -template -mockup -cartoon -people -person -human -man -woman -child -children -students -teacher -face -portrait"
RELAXED_NEGATIVE = " -icon -icons -logo -svg -clipart -cartoon -people -person"   # lượt cứu cánh

# Contexts cho ảnh objects/concepts (KHÔNG có người)
OBJECT_CONTEXTS = [
    "technology concept",
    "software development",
    "programming tools",
    "computer system",
    "digital technology",
    "application interface",
    "software architecture",
    "technical illustration",
    "concept visualization",
    "system diagram"
]

# Hiển thị caption (heading) lên ảnh
TEXT_OVERLAY = True

# Ưu tiên các khóa URL ảnh thường gặp từ nhiều provider
CAND_URL_KEYS = ("image_url", "image", "contentUrl", "thumbnail", "thumb", "src", "url")


# ======================================================================
# Utils
# ======================================================================

def _safe_makedirs(p: str):
    os.makedirs(p, exist_ok=True)


def _best_image_url(cand: dict) -> str:
    for k in CAND_URL_KEYS:
        u = cand.get(k)
        if isinstance(u, str) and u.startswith("http"):
            return u
    return ""


# ======================================================================
# ĐỌC VĂN BẢN
# ======================================================================

def _read_pdf(path: str) -> str:
    txt = []
    with open(path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            t = (page.extract_text() or "").strip()
            if t:
                txt.append(t)
    return "\n\n".join(txt)


def _read_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def _read_pptx(path: str) -> str:
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        buf = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                t = shape.text.strip()
                if t:
                    buf.append(t)
        if buf:
            chunks.append("\n".join(buf))
    return "\n\n".join(chunks)


def extract_text_from_file(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf": return _read_pdf(path)
    if ext == ".docx": return _read_docx(path)
    if ext in (".pptx", ".ppt"): return _read_pptx(path)
    if ext == ".txt":
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    raise ValueError(f"Unsupported format: {ext}")


def clean_text(s: str) -> str:
    s = re.sub(r'[ \t]+\n', '\n', s)
    s = re.sub(r'\n{3,}', '\n\n', s)
    return s.strip()


# ======================================================================
# CHIA KHỐI NỘI DUNG
# ======================================================================

def chunk_text(s: str, max_chars: int = MAX_CHARS_PER_CHUNK) -> List[str]:
    sentences = re.split(r'(?<=[\.!?])\s+', s)
    chunks, cur = [], ""
    for sent in sentences:
        if len(cur) + len(sent) + 1 <= max_chars:
            cur = (cur + " " + sent).strip()
        else:
            if cur: chunks.append(cur)
            if len(sent) <= max_chars:
                cur = sent
            else:
                wraps = textwrap.wrap(sent, width=90)
                for w in wraps:
                    if len(w) > max_chars:
                        for sub in textwrap.wrap(w, width=70):
                            chunks.append(sub)
                    else:
                        chunks.append(w)
                cur = ""
    if cur: chunks.append(cur)
    return [c.strip() for c in chunks if len(c.strip()) >= 20][:MAX_SLIDES]


def estimate_duration_sec(text: str, wpm: int = WPM) -> float:
    words = max(1, len(text.split()))
    return round(words / max(80, wpm) * 60.0, 2)


def _improve_text_for_tts(text: str) -> str:
    text = re.sub(r'[^\w\s.,!?;:()\-]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text + ('.' if not text.endswith(('.', '!', '?')) else '')


# ======================================================================
# TỪ KHOÁ/CHỦ ĐỀ & VẼ CHỮ
# ======================================================================

def _extract_keywords(text: str, top_n: int = 5, full_content: str = "") -> list:
    """Cải thiện extract keywords: ưu tiên từ sau headings và từ lặp lại nhiều lần trong file"""
    # 1) Extract từ headings (từ đằng sau các đề mục lớn) - CẢI THIỆN
    heading_keywords = []
    if full_content:
        # Tìm các dòng bắt đầu bằng số hoặc chữ hoa (headings)
        lines = full_content.split('\n')
        for line in lines:
            line_stripped = line.strip()
            # Heading patterns: "1. Text", "1.1 Text", "Chapter X", "## Title", etc.
            is_heading = (
                re.match(r'^[\d#\-\s.]+\s+[A-Z]', line_stripped) or  # "1. Title", "1.1. Title"
                (re.match(r'^[A-Z][A-Za-z\s]{5,}', line_stripped[:80]) and len(line_stripped) < 100) or  # Title case, short
                re.match(r'^Chapter\s+\d+', line_stripped, re.I) or  # "Chapter 1"
                re.match(r'^#{1,3}\s+[A-Z]', line_stripped)  # Markdown headings
            )
            if is_heading:
                # Lấy cụm danh từ đằng sau đề mục - CẢI THIỆN để lấy cụm từ tốt hơn
                clean_heading = re.sub(r'^[\d#\-\s.Chapter]+', '', line_stripped, flags=re.I)
                clean_heading = re.sub(r'[^\w\s]', ' ', clean_heading)
                
                # 1) Ưu tiên cụm danh từ (noun phrases): từ viết hoa liền nhau hoặc từ dài
                words = clean_heading.split()
                noun_phrases = []
                current_phrase = []
                
                for i, w in enumerate(words):
                    # Nếu từ viết hoa hoặc từ dài (>= 5), có thể là danh từ quan trọng
                    if (w[0].isupper() or len(w) >= 5) and len(w) >= 4 and w.lower() not in STOPWORDS:
                        current_phrase.append(w.lower())
                    else:
                        # Kết thúc cụm danh từ hiện tại
                        if len(current_phrase) >= 2:
                            noun_phrases.append(" ".join(current_phrase))
                        elif len(current_phrase) == 1 and len(current_phrase[0]) >= 5:
                            heading_keywords.append(current_phrase[0])
                        current_phrase = []
                        # Vẫn thêm từ dài đơn lẻ
                        if len(w) >= 5 and w.lower() not in STOPWORDS:
                            heading_keywords.append(w.lower())
                
                # Thêm cụm cuối cùng
                if len(current_phrase) >= 2:
                    noun_phrases.append(" ".join(current_phrase))
                elif len(current_phrase) == 1 and len(current_phrase[0]) >= 5:
                    heading_keywords.append(current_phrase[0])
                
                # Thêm cụm danh từ vào keywords
                heading_keywords.extend(noun_phrases[:3])  # Top 3 cụm từ
                
                # 2) Lấy từ đơn quan trọng (backup)
                single_words = [w.lower() for w in words 
                               if len(w) >= 5 and w.lower() not in STOPWORDS and w.lower() not in heading_keywords]
                heading_keywords.extend(single_words[:2])  # Thêm 2 từ đơn
                
                if noun_phrases or heading_keywords:
                    print(f"[KEYWORD] Extracted from heading '{line_stripped[:60]}': phrases={noun_phrases[:3]}, words={[w for w in heading_keywords[-5:] if w not in noun_phrases]}")
    
    # 2) Tìm từ lặp lại nhiều lần trong toàn bộ file (nếu có full_content) hoặc chỉ trong text hiện tại
    search_text = full_content if full_content and len(full_content) > len(text) * 2 else text
    clean = re.sub(r"[^A-Za-z0-9\s\-]", " ", search_text)
    words = [w for w in clean.split() if len(w) >= 4 and w.lower() not in STOPWORDS]
    
    # Tính frequency - từ xuất hiện nhiều lần sẽ có điểm cao hơn
    freq: Dict[str, int] = {}
    for w in words:
        k = w.lower()
        freq[k] = freq.get(k, 0) + 1
    
    # Ưu tiên: heading keywords > frequent words (frequency cao và từ dài)
    # Heading keywords được boost điểm
    for kw in heading_keywords:
        if kw in freq:
            freq[kw] = freq[kw] * 2  # Boost heading keywords
    
    # Sắp xếp: frequency * length (ưu tiên từ dài và xuất hiện nhiều)
    ranked = sorted(freq.items(), key=lambda x: (x[1] * len(x[0]), x[1], len(x[0])), reverse=True)
    
    # Lấy top keywords - ưu tiên heading keywords trước
    keywords = []
    added = set()
    
    # Thêm heading keywords trước
    for kw in heading_keywords:
        if kw in freq and kw not in added:
            keywords.append(kw)
            added.add(kw)
    
    # Thêm các từ frequent khác
    for w, count in ranked:
        if w not in added and count >= 2:  # Chỉ lấy từ xuất hiện ít nhất 2 lần
            keywords.append(w)
            added.add(w)
            if len(keywords) >= top_n * 2:  # Lấy nhiều hơn để filter
                break
    
    # Thêm bigrams từ text hiện tại nếu cần
    if len(keywords) < top_n:
        bigrams = []
        words_lower = [w.lower() for w in clean.split() if len(w) >= 3]
        for i in range(len(words_lower) - 1):
            bigram = f"{words_lower[i]} {words_lower[i+1]}"
            if (len(bigram.split()[0]) >= 4 and len(bigram.split()[1]) >= 3 and
                bigram not in added):
                bigrams.append(bigram)
        keywords.extend(bigrams[:2])
    
    return keywords[:top_n]  # Trả về top_n keywords tốt nhất


def _looks_thematic(meta: str, query: str = "") -> bool:
    s = ((meta or "") + " " + (query or "")).lower()
    # Nới lỏng: chỉ reject nếu banned keyword là từ độc lập (word boundary), không phải substring
    # Ví dụ: "classroom" không bị reject vì có "room" trong "classroom"
    for bad in BANNED_KEYWORDS:
        # Kiểm tra nếu banned keyword xuất hiện như một từ độc lập hoặc ở đầu/cuối
        pattern = r'\b' + re.escape(bad) + r'\b'
        if re.search(pattern, s):
            # Kiểm tra thêm: chỉ reject nếu không phải là từ trong từ khác
            # Ví dụ: "classroom" chứa "room" nhưng không nên reject
            words_in_s = set(s.split())
            if bad in words_in_s or s.startswith(bad) or s.endswith(bad):
                print(f"[REJECT] Banned keyword '{bad}' in meta: {meta[:100]}")
                return False
    return True


def _extract_main_noun_from_heading(heading: str) -> str:
    clean = re.sub(r'^\d+\.\s*', '', heading.strip())
    clean = re.sub(r'[^\w\s]', ' ', clean)
    words = clean.split()
    proper_nouns = [w for w in words if w[:1].isupper() and len(w) > 2]
    if proper_nouns:
        return proper_nouns[0]
    long_words = [w for w in words if len(w) > 4]
    if long_words:
        return max(long_words, key=len)
    return "topic"  # Thay "education" bằng general để linh hoạt


def _overlay_caption(img: Image.Image, caption: str) -> Image.Image:
    if not caption:
        return img
    img = img.copy()
    draw = ImageDraw.Draw(img, "RGBA")
    W, H = img.size
    pad = 24
    try:
        font = ImageFont.truetype("arial.ttf", 36)
    except Exception:
        font = ImageFont.load_default()
    try:
        bbox = draw.textbbox((0, 0), caption, font=font)
        tw, th = bbox[2], bbox[3]
    except Exception:
        tw, th = (len(caption) * 9, 18)
    box_h = th + pad
    draw.rectangle([0, H - box_h, W, H], fill=(0, 0, 0, 140))
    draw.text((pad // 2, H - box_h + (pad // 3)), caption, fill=(255, 255, 255, 230), font=font)
    return img


# ======================================================================
# ẢNH DỰ PHÒNG (gradient – không icon/khối/ chữ)
# ======================================================================

def _create_fallback_image(tmpdir: str, variation: int, caption: Optional[str] = None) -> str:
    """Fallback: nền gradient dịu, KHÔNG chữ/khối. Variation khác màu để đa dạng."""
    W, H = VIDEO_WIDTH, VIDEO_HEIGHT

    # Thay đổi màu theo variation để không lặp 1 hình trắng
    colors = [
        ([210, 220, 235], [245, 248, 255]),  # Xanh nhạt - trắng
        ([220, 240, 255], [240, 255, 250]),  # Xanh dương - xanh nhạt
        ([255, 240, 220], [255, 250, 240]),  # Cam nhạt - trắng
        ([230, 255, 230], [245, 255, 245])   # Xanh lá nhạt - trắng
    ]
    top, bot = colors[variation % len(colors)]
    top = np.array(top, dtype=np.float32)
    bot = np.array(bot, dtype=np.float32)

    t = np.linspace(0.0, 1.0, H, dtype=np.float32).reshape(H, 1, 1)
    grad_col = top.reshape(1, 1, 3) * (1.0 - t) + bot.reshape(1, 1, 3) * t
    img_arr = np.repeat(grad_col, W, axis=1).astype(np.uint8)

    img = Image.fromarray(img_arr, mode="RGB")
    if caption:
        caption_text = " ".join(caption.split())[:70]
        img = _overlay_caption(img, caption_text)
    path = os.path.join(tmpdir, f"fallback_{variation}.png")
    img.save(path, quality=92)
    return path

# ======================================================================
# ẢNH: helper chất lượng/đa dạng
# ======================================================================

def _fit_to_frame(img: Image.Image, w: int = VIDEO_WIDTH, h: int = VIDEO_HEIGHT) -> Image.Image:
    return ImageOps.fit(img, (w, h), method=Image.BICUBIC, centering=(0.5, 0.5))

def _tokenize(s: str) -> List[str]:
    return [w.lower() for w in re.findall(r"[A-Za-z0-9]+", s or "") if len(w) >= 3]

def _score_meta(meta: str, query: str, heading: str, extra_kws: List[str]) -> float:
    tokens = set(_tokenize(meta))
    base = set(_tokenize(heading)) | set(map(str.lower, extra_kws)) | set(_tokenize(query))
    if not tokens:
        return 0.0
    overlap = tokens & base
    score = len(overlap) * 2.0 + min(5.0, len(tokens) / 10.0)
    # Boost nếu khớp anchor giáo dục (không bắt buộc)
    meta_lower = meta.lower()
    anchor_matches = sum(1 for anchor in TOPIC_ANCHORS if anchor in meta_lower)
    score += anchor_matches * 1.5
    return score

def _entropy(img: Image.Image) -> float:
    g = img.convert("L").resize((64, 64), Image.BICUBIC)
    hist = g.histogram()
    total = sum(hist)
    ent = 0.0
    for c in hist:
        if c:
            p = c / total
            ent -= p * math.log2(p)
    return ent

def _edge_score(img: Image.Image) -> float:
    g = img.convert("L").resize((256, 256), Image.BICUBIC)
    e = g.filter(ImageFilter.FIND_EDGES)
    return float(np.array(e, dtype=np.float32).mean())

def _is_text_only_image(img: Image.Image) -> bool:
    ent = _entropy(img)
    uniq = len(set(img.convert("P", palette=Image.ADAPTIVE, colors=64).getdata()))
    return ent < 3.0 and uniq < 32

def _is_icon_or_flat_illustration(img: Image.Image) -> bool:
    uniq = len(set(img.convert("P", palette=Image.ADAPTIVE, colors=64).getdata()))
    ent = _entropy(img)
    edge = _edge_score(img)
    # Nới ngưỡng để ít loại nhầm ảnh thực tế hơn
    return (uniq < 16 and ent < 2.0) and (edge < 4.0)

def _ahash(img: Image.Image, hash_size: int = HASH_SIZE) -> int:
    g = img.convert("L").resize((hash_size, hash_size), Image.BICUBIC)
    pixels = list(g.getdata())
    avg = sum(pixels) / len(pixels)
    bits = 0
    for p in pixels:
        bits = (bits << 1) | (1 if p > avg else 0)
    return bits

def _hamming(a: int, b: int) -> int:
    x = a ^ b
    return x.bit_count() if hasattr(int, "bit_count") else bin(x).count("1")

def _quality_ok(img: Image.Image) -> bool:
    return img.width >= MIN_IMG_W and img.height >= MIN_IMG_H

def _host_from_url(u: str) -> str:
    try:
        return re.sub(r"^www\.", "", re.findall(r"://([^/]+)", u)[0].lower())
    except Exception:
        return ""


# ======================================================================
# TÌM ẢNH THEO SLIDE (liên quan & đa dạng)
# ======================================================================

def _find_images_for_slide(text: str, slide_num: int, tmpdir: str, full_content: str = "") -> List[str]:
    image_paths: List[str] = []
    t0 = time.monotonic()
    try:
        # 1) Heading theo số slide
        heading = ""
        lines = full_content.split('\n')
        current_slide = f"{slide_num}."
        for line in lines:
            if line.strip().startswith(current_slide) or re.match(rf'^\s*{slide_num}\s*(?:[.\-)]\s+)', line.strip()):
                heading = line.strip()
                break
        if not heading:
            heading = (text.split("\n")[0] or "").strip()

        main_kw = _extract_main_noun_from_heading(heading) or "topic"
        extra_kws = _extract_keywords(text, top_n=5, full_content=full_content)  # Lấy keywords từ headings và từ lặp lại

        # 2) Truy vấn "ảnh thật" (không icon/clipart/mockup) - Cải thiện để khớp nội dung chính xác hơn
        # Tạo queries dựa trên keywords thực tế từ nội dung và bối cảnh giáo dục
        base1 = " ".join([main_kw] + extra_kws[:3]).strip() if extra_kws else main_kw
        base2 = (heading.split("\n")[0][:80] if heading else main_kw).strip()

        queries: List[str] = []

        def _add_query(q: str):
            if q:
                queries.append(q + NEGATIVE_PROMPT)

        # Query chính: TẬP TRUNG VÀO OBJECTS/CONCEPTS, KHÔNG CÓ NGƯỜI
        # Ưu tiên keywords từ headings và từ lặp lại nhiều lần
        focus_keywords = [main_kw] + extra_kws[:5]  # Top 5 keywords quan trọng nhất
        
        # Query 1: Từ khóa chính trực tiếp (objects/concepts)
        for kw in focus_keywords[:3]:  # Top 3 keywords
            if kw and len(kw) >= 4 and kw.lower() not in STOPWORDS:
                # Query trực tiếp về keyword (object/concept)
                _add_query(kw)
                # Query với technology/software context (không có người)
                _add_query(f"{kw} technology")
                _add_query(f"{kw} software")
                _add_query(f"{kw} system")
        
        # Query 2: Cụm từ từ headings (nếu có cụm từ 2+ từ)
        for kw in focus_keywords:
            if " " in kw and len(kw) >= 8:  # Cụm từ dài
                _add_query(kw)
                _add_query(f"{kw} concept")
                _add_query(f"{kw} illustration")
        
        # Query 3: Kết hợp keywords chính với object contexts (KHÔNG có người)
        for kw in focus_keywords[:4]:
            if kw and len(kw) >= 4:
                for ctx in OBJECT_CONTEXTS[:4]:  # Top 4 contexts
                    _add_query(f"{kw} {ctx}")

        queries = list(dict.fromkeys(queries))[:MAX_QUERIES_PER_SLIDE]
        print(f"[QUERY] Slide {slide_num} queries: {queries}")

        candidates, seen_urls = [], set()

        def _try_download(url: str, used_query: str, idx_for_caption: int,
                          min_w: int = MIN_IMG_W, min_h: int = MIN_IMG_H,
                          skip_icon_check: bool = False) -> None:
            if not url or url in seen_urls:
                return
            try:
                r = requests.get(
                    url,
                    headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"},
                    timeout=PER_IMAGE_HTTP_TIMEOUT
                )
                if r.status_code != 200:
                    print(f"[REJECT] Download fail status {r.status_code} for url: {url}")
                    return
                ct = (r.headers.get("content-type") or "").lower()
                if ("image" not in ct) and (not url.lower().endswith((".jpg", ".jpeg", ".png", ".webp", ".bmp"))):
                    print(f"[REJECT] Invalid content-type for url: {url}")
                    return

                img_raw = Image.open(BytesIO(r.content)).convert("RGB")
                if img_raw.width < min_w or img_raw.height < min_h:
                    print(f"[REJECT] Size too small {img_raw.width}x{img_raw.height} for url: {url}")
                    return
                if _is_text_only_image(img_raw):
                    print(f"[REJECT] Text-only image for url: {url}")
                    return
                if (not skip_icon_check) and _is_icon_or_flat_illustration(img_raw):
                    print(f"[REJECT] Icon/flat illustration for url: {url}")
                    return

                ah = _ahash(img_raw)
                host = _host_from_url(url)
                img_fit = _fit_to_frame(img_raw, VIDEO_WIDTH, VIDEO_HEIGHT)
                if TEXT_OVERLAY:
                    caption = heading if idx_for_caption == 0 else f"{main_kw.title()}"
                    img_fit = _overlay_caption(img_fit, caption)
                p = os.path.join(tmpdir, f"img_{slide_num}_{len(candidates)}.jpg")
                img_fit.save(p, quality=92)

                score = _score_meta(url, used_query, heading, extra_kws)
                candidates.append({
                    "path": p, "score": score, "host": host, "ahash": ah,
                    "w": img_raw.width, "h": img_raw.height
                })
                seen_urls.add(url)
            except Exception as e:
                print(f"[REJECT] Download exception {e} for url: {url}")

        # 3) Thu thập ứng viên (có time-budget)
        for q in queries:
            if time.monotonic() - t0 > PER_SLIDE_TIME_BUDGET:
                break
            try:
                images = search_images(q, count=MAX_FETCH_PER_QUERY, orientation="any")
            except Exception as e:
                print(f"[ERROR] search_images fail for query {q}: {e}")
                images = []
            for cand in images:
                if time.monotonic() - t0 > PER_SLIDE_TIME_BUDGET:
                    break
                url = _best_image_url(cand)
                meta = " ".join([str(cand.get('title','')), str(cand.get('alt','')),
                                 str(cand.get('source','')), url])
                if not _looks_thematic(meta, q):
                    print(f"[REJECT] Thematic fail for meta: {meta[:100]} url: {url}")
                    continue
                _try_download(url, q, idx_for_caption=len(image_paths))

        # 4) Không có ứng viên → lượt cứu cánh (nới ràng buộc) - Tăng retry để tránh fallback
        retry_count = 0
        while len(candidates) < FINAL_IMAGES_PER_SLIDE and retry_count < 3:  # Retry 3 lần
            retry_count += 1
            rescue_queries = []
            for base in [base1, base2]:
                if base:
                    rescue_queries.append(f"real education photo of {base}" + RELAXED_NEGATIVE)
            for kw in focus_keywords:
                if kw:
                    for ctx in OBJECT_CONTEXTS[:5]:
                        rescue_queries.append(f"{kw} {ctx}" + RELAXED_NEGATIVE)
            # Thêm truy vấn chung
            rescue_queries.append(f"education classroom learning {main_kw}" + RELAXED_NEGATIVE)
            rescue_queries = list(set(rescue_queries))[:max(6, MAX_QUERIES_PER_SLIDE)]
            print(f"[RESCUE] Retry {retry_count} queries: {rescue_queries}")

            for q in rescue_queries:
                if time.monotonic() - t0 > PER_SLIDE_TIME_BUDGET * (1.5 + retry_count * 0.5):
                    break
                try:
                    images = search_images(q, count=MAX_FETCH_PER_QUERY * 2, orientation="any")
                except Exception as e:
                    print(f"[ERROR] Rescue search fail for {q}: {e}")
                    images = []
                for cand in images:
                    if time.monotonic() - t0 > PER_SLIDE_TIME_BUDGET * (1.5 + retry_count * 0.5):
                        break
                    url = _best_image_url(cand)
                    meta = " ".join([str(cand.get('title','')), str(cand.get('alt','')),
                                     str(cand.get('source','')), url])
                    if not _looks_thematic(meta, q):
                        print(f"[REJECT] Rescue thematic fail for meta: {meta[:100]} url: {url}")
                        continue
                    _try_download(url, q, idx_for_caption=len(image_paths),
                                  min_w=800, min_h=500, skip_icon_check=True)  # Nới nhưng giữ chất, skip icon check

        # 5) Vẫn trống → fallback (bây giờ hiếm, chỉ nếu ==0)
        if len(candidates) == 0:
            print(f"[FALLBACK] No candidates after all retries for slide {slide_num}")
            fallback_caption = (heading or main_kw or "education lesson")
            for i in range(FINAL_IMAGES_PER_SLIDE):
                image_paths.append(_create_fallback_image(tmpdir, slide_num * 10 + i, fallback_caption))
            return image_paths

        # 6) Chọn ảnh: điểm cao + domain đa dạng + hash khác nhau
        def pick(require_domain_diversity=True):
            chosen, hosts, hashes = [], set(), []
            sorted_c = sorted(candidates, key=lambda c: (c["score"], c["w"] * c["h"]), reverse=True)
            for c in sorted_c:
                if len(chosen) >= FINAL_IMAGES_PER_SLIDE:
                    break
                if require_domain_diversity and c["host"] in hosts:
                    continue
                if any(_hamming(c["ahash"], h) < MIN_HASH_DISTANCE for h in hashes):
                    continue
                chosen.append(c["path"]); hashes.append(c["ahash"]); hosts.add(c["host"])
            return chosen

        picks = pick(True)
        if len(picks) < FINAL_IMAGES_PER_SLIDE:
            sorted_c = sorted(candidates, key=lambda c: (c["score"], c["w"] * c["h"]), reverse=True)
            for c in sorted_c:
                if len(picks) >= FINAL_IMAGES_PER_SLIDE: break
                if c["path"] in picks: continue
                if any(_hamming(c["ahash"], _c["ahash"]) < MIN_HASH_DISTANCE
                       for _c in candidates if _c.get("path") in picks): continue
                picks.append(c["path"])
        if len(picks) < FINAL_IMAGES_PER_SLIDE:
            rest = [c["path"] for c in candidates if c["path"] not in picks]
            picks += rest[: FINAL_IMAGES_PER_SLIDE - len(picks)]

        # Randomize số ảnh: 2-3 ảnh mỗi slide
        num_images = random.randint(2, 3)
        image_paths = picks[:num_images]
        print(f"[IMG] slide {slide_num}: candidates={len(candidates)} picked={len(image_paths)} (requested={num_images})")
        return image_paths

    except Exception as e:
        print(f"[IMAGE] Lỗi tổng: {e}")
        fb = _create_fallback_image(tmpdir, slide_num)
        return [fb] * FINAL_IMAGES_PER_SLIDE


# ======================================================================
# XỬ LÝ SLIDE SONG SONG
# ======================================================================

def _process_slide_parallel(
    slide_idx: int,
    text: str,
    total_slides: int,
    width: int,
    height: int,
    title: str,
    tmpdir: str,
    variation: int,
    full_content: str
) -> Tuple[int, List[str], str, float]:
    try:
        frame_paths = _find_images_for_slide(text, slide_idx, tmpdir, full_content)
        wav_path = os.path.join(tmpdir, f"tts_{slide_idx:03d}.wav")
        duration = synth_voice_only_human(text, wav_path)
        audio_path = wav_path if os.path.exists(wav_path) else ""
        if audio_path:
            try:
                # Validate audio file before returning, otherwise fallback to silent
                with contextlib.closing(AudioFileClip(audio_path)) as test_clip:
                    test_clip.duration  # force load metadata
            except Exception as audio_err:
                print(f"[AUDIO] Invalid audio for slide {slide_idx}: {audio_err}, using silent fallback")
                audio_path = ""
        return slide_idx, frame_paths, audio_path, duration
    except Exception as e:
        print(f"[Slide {slide_idx} Error] {e}")
        caption = (text.split("\n")[0] or "education lesson")
        fallback = _create_fallback_image(tmpdir, slide_idx, caption)
        return slide_idx, [fallback] * FINAL_IMAGES_PER_SLIDE, "", estimate_duration_sec(text)


# ======================================================================
# TTS AN TOÀN — NHANH
# ======================================================================

def synth_voice_only_human(text: str, out_path: str) -> float:
    improved = _improve_text_for_tts(text)
    duration = estimate_duration_sec(improved)

    # 1) Nếu yêu cầu giọng người thật nhưng không khả dụng → báo lỗi
    human_available = is_human_voice_available()
    if REQUIRE_HUMAN_VOICE_ONLY and not human_available:
        raise RuntimeError("Human voice is required but not available. Please configure ElevenLabs API key.")

    # 2) ElevenLabs nếu có
    if human_available:
        try:
            result = synthesize_human_voice(improved, output_path=out_path)
            if result.get("success") and os.path.exists(out_path):
                print(f"[TTS] ElevenLabs OK: {out_path}")
                return result.get("duration") or duration
            print("[TTS] ElevenLabs response missing success flag, will fallback")
        except Exception as e:
            print(f"[TTS] ElevenLabs failed: {e}")
            if REQUIRE_HUMAN_VOICE_ONLY:
                raise RuntimeError("Human voice synthesis failed and REQUIRE_HUMAN_VOICE is enabled") from e

    if REQUIRE_HUMAN_VOICE_ONLY:
        raise RuntimeError("Human voice synthesis failed and REQUIRE_HUMAN_VOICE is enabled")

    # 3) Fallback: gTTS (giọng máy) nếu được phép
    try:
        from gtts import gTTS
        # Tạo file MP3 tạm với tên rõ ràng
        base_name = os.path.splitext(out_path)[0]
        mp3_path = base_name + "_gtts.mp3"
        wav_path = out_path if out_path.lower().endswith(".wav") else base_name + ".wav"
        
        tts = gTTS(text=improved[:500], lang='en', slow=False)
        tts.save(mp3_path)
        print(f"[TTS] gTTS MP3 created: {mp3_path}")
        
        # Transcode MP3 -> WAV
        try:
            aclip = AudioFileClip(mp3_path)
            if aclip.duration and aclip.duration > 0:
                aclip.write_audiofile(wav_path, fps=44100, logger=None)
                print(f"[TTS] gTTS transcoded to WAV: {wav_path} (duration={aclip.duration:.2f}s)")
                aclip.close()
            else:
                raise ValueError("gTTS audio has invalid duration")
        finally:
            # Cleanup MP3 file
            if os.path.exists(mp3_path):
                os.remove(mp3_path)
        
        # Đảm bảo out_path trỏ đến file WAV đúng
        if wav_path != out_path:
            if os.path.exists(wav_path):
                # Nếu out_path đã tồn tại, xóa nó
                if os.path.exists(out_path):
                    os.remove(out_path)
                # Copy hoặc rename WAV to out_path
                shutil.copy2(wav_path, out_path)
                # Xóa file WAV tạm nếu khác với out_path
                if wav_path != out_path and os.path.exists(wav_path):
                    os.remove(wav_path)
        
        return duration
    except Exception as e:
        print(f"[TTS] gTTS failed: {e}")
        # 4) Cuối cùng mới silent nếu tất cả fail
        try:
            samples = int(44100 * duration)
            audio_array = np.zeros((samples, 2), dtype=np.float32)
            clip = AudioArrayClip(audio_array, fps=44100)
            clip.write_audiofile(out_path, fps=44100, logger=None)
            print(f"[TTS] Silent fallback: {out_path}")
            return duration
        except Exception as e2:
            print(f"[TTS] Silent fallback failed: {e2}")
            return duration


def _close_clip(clip):
    try:
        if clip is not None:
            clip.close()
            if hasattr(clip, "audio") and clip.audio:
                clip.audio.close()
    except Exception:
        pass


def _make_silent_audio(duration: float) -> AudioArrayClip:
    samples = int(44100 * duration)
    return AudioArrayClip(np.zeros((samples, 2), dtype=np.float32), fps=44100)


def write_vtt(cues: List[Tuple[float, float, str]], path: str):
    def fmt(t):
        ms = int((t - int(t)) * 1000)
        s = int(t)
        h, m, s = s // 3600, (s % 3600) // 60, s % 60
        return f"{h:02d}:{m:02d}:{s:02d}.{ms:03d}"
    lines = ["WEBVTT", ""]
    for start, end, text in cues:
        lines.append(f"{fmt(start)} --> {fmt(end)}")
        lines.append(text.strip())
        lines.append("")
    _safe_makedirs(os.path.dirname(path))
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


class TempWorkspace:
    def __init__(self):
        self.tmpdir = tempfile.mkdtemp(prefix="doc2vid_")
    def __enter__(self):
        return self
    def __exit__(self, *args):
        shutil.rmtree(self.tmpdir, ignore_errors=True)
        gc.collect()


# ======================================================================
# MAIN
# ======================================================================

def make_video_from_file(
    src_path: str,
    out_dir: str,
    title: str = None,
    width: int = VIDEO_WIDTH,
    height: int = VIDEO_HEIGHT,
    wpm: int = WPM
) -> Dict[str, str]:
    _safe_makedirs(out_dir)

    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)

    if len(raw) > MAX_TEXT_LENGTH:
        print(f"[WARN] File lớn ({len(raw)} ký tự), cắt xuống {MAX_TEXT_LENGTH}...")
        raw = raw[:MAX_TEXT_LENGTH]

    title = title or os.path.splitext(os.path.basename(src_path))[0]
    chunks = chunk_text(raw)[:MAX_SLIDES]
    if not chunks:
        raise ValueError("Không có nội dung.")

    # Giới hạn thời lượng tổng
    estimated_total = sum(estimate_duration_sec(chunk) for chunk in chunks)
    if estimated_total > MAX_DURATION_SEC:
        max_slides_by_time = int(MAX_DURATION_SEC / (estimated_total / len(chunks)))
        max_slides_by_time = max(1, min(max_slides_by_time, MAX_SLIDES))
        chunks = chunks[:max_slides_by_time]
        print(f"[OPTIMIZE] Giới hạn {max_slides_by_time} slides (<{MAX_DURATION_SEC}s)")

    # Đảm bảo mỗi video có VTT file riêng, không trùng lặp
    # Input: {uuid}_{filename}.pdf -> Output: {uuid}_{filename}_video.mp4 và {uuid}_{filename}.vtt
    out_base = os.path.splitext(os.path.basename(src_path))[0]
    # out_base đã có UUID rồi từ api_doc_to_video, nên không cần thêm UUID nữa
    out_video = os.path.join(out_dir, f"{out_base}_video.mp4")
    out_vtt = os.path.join(out_dir, f"{out_base}.vtt")
    print(f"[VTT] Creating unique VTT file: {out_vtt} (for video: {out_video})")

    clips: List[ImageClip] = []
    timeline: List[Tuple[float, float, str]] = []
    cur_time = 0.0

    with TempWorkspace() as ws:
        tmpdir = ws.tmpdir
        results: Dict[int, Tuple[int, List[str], str, float]] = {}
        # Lưu mapping index -> chunk để tránh lỗi
        chunk_map = {i + 1: chunks[i] for i in range(len(chunks))}

        # Xử lý ảnh + TTS song song
        with ThreadPoolExecutor(max_workers=min(6, os.cpu_count() or 1)) as exec:
            futures = {
                exec.submit(_process_slide_parallel, i + 1, chunk, len(chunks), width, height, title, tmpdir, i, raw): i + 1
                for i, chunk in enumerate(chunks)
            }
            for f in as_completed(futures):
                i = futures[f]
                try:
                    results[i] = f.result()
                except Exception as e:
                    print(f"[Future Error] Slide {i}: {e}")
                    fallback_caption = chunk_map.get(i, "") or "education lesson"
                    fallback_img = _create_fallback_image(tmpdir, i, fallback_caption)
                    results[i] = (i, [fallback_img] * FINAL_IMAGES_PER_SLIDE, "", 5.0)

        # Dựng video
        for i in sorted(results.keys()):
            idx, image_paths, audio_path, dur = results[i]
            # Lấy text chunk cho slide này để dùng cho audio fallback
            # KHÔNG đặt tên là chunk_text vì trùng với hàm chunk_text()
            current_chunk_text = chunk_map.get(i, "")
            if not isinstance(image_paths, list) or not image_paths:
                fallback_img = _create_fallback_image(tmpdir, idx, current_chunk_text or "education lesson")
                image_paths = [fallback_img] * FINAL_IMAGES_PER_SLIDE
            num_images = max(1, len(image_paths))
            duration_per_image = max(1.0, dur / num_images)

            subclips: List[ImageClip] = []
            for img_idx, img_path in enumerate(image_paths):
                try:
                    clip = ImageClip(img_path).set_duration(duration_per_image).resize((width, height))
                    subclips.append(clip)
                except Exception as e:
                    print(f"[SUBCLIP] Lỗi ảnh {img_idx}: {e}")
                    fallback_path = _create_fallback_image(tmpdir, idx * 10 + img_idx, current_chunk_text or "education lesson")
                    clip = ImageClip(fallback_path).set_duration(duration_per_image).resize((width, height))
                    subclips.append(clip)

            if subclips:
                # Tạo slide clip với tất cả ảnh, mỗi ảnh có duration bằng nhau
                slide_clip = concatenate_videoclips(subclips, method="compose").set_fps(ENCODE_FPS)
                print(f"[SUBCLIPS] Created slide {idx} with {len(subclips)} images, total duration={slide_clip.duration:.2f}s")
            else:
                fb = _create_fallback_image(tmpdir, idx, current_chunk_text or "education lesson")
                slide_clip = ImageClip(fb).set_duration(dur).resize((width, height)).set_fps(ENCODE_FPS)

            # Gắn audio - ưu tiên human voice, fallback về gTTS nếu không có
            # QUAN TRỌNG: Điều chỉnh duration của từng ảnh thay vì cắt video
            if audio_path and os.path.exists(audio_path):
                try:
                    # Kiểm tra file size trước - file quá nhỏ có thể là empty
                    file_size = os.path.getsize(audio_path)
                    if file_size < 1000:  # < 1KB có thể là file rỗng
                        print(f"[AUDIO] File too small ({file_size} bytes), treating as empty")
                        raise ValueError("Audio file too small")
                    
                    aclip = AudioFileClip(audio_path)
                    audio_duration = aclip.duration
                    
                    # Kiểm tra duration hợp lệ
                    if audio_duration is None or audio_duration <= 0:
                        print(f"[AUDIO] Invalid audio duration: {audio_duration}")
                        raise ValueError("Invalid audio duration")
                    
                    video_duration = slide_clip.duration
                    
                    # Điều chỉnh video duration để khớp với audio - thay vì cắt, tái tạo với duration mới
                    if video_duration != audio_duration:
                        # Tính lại duration cho mỗi ảnh
                        new_duration_per_image = max(1.0, audio_duration / num_images)
                        # Tái tạo subclips với duration mới
                        adjusted_subclips = []
                        for img_idx, img_path in enumerate(image_paths):
                            try:
                                clip = ImageClip(img_path).set_duration(new_duration_per_image).resize((width, height))
                                adjusted_subclips.append(clip)
                            except Exception:
                                # Fallback nếu có lỗi
                                fallback_path = _create_fallback_image(tmpdir, idx * 10 + img_idx, current_chunk_text or "education lesson")
                                clip = ImageClip(fallback_path).set_duration(new_duration_per_image).resize((width, height))
                                adjusted_subclips.append(clip)
                        
                        if adjusted_subclips:
                            slide_clip = concatenate_videoclips(adjusted_subclips, method="compose").set_fps(ENCODE_FPS)
                            video_duration = slide_clip.duration
                            print(f"[DURATION FIX] Adjusted {num_images} images to match audio: {video_duration:.2f}s (per image: {new_duration_per_image:.2f}s)")
                    elif video_duration < audio_duration:
                        # Video ngắn hơn audio: cắt audio cho vừa với video
                        aclip = aclip.subclip(0, video_duration)
                        audio_duration = video_duration
                        print(f"[DURATION FIX] Audio cut to match video: {audio_duration:.2f}s")
                    
                    # Đảm bảo audio duration khớp với video duration (backup check)
                    if audio_duration < video_duration:
                        # Audio ngắn hơn: loop để đủ duration
                        if audio_duration > 0:
                            loops_needed = int(np.ceil(video_duration / audio_duration)) + 1
                            audio_clips = [aclip] * loops_needed
                            aclip_looped = concatenate_audioclips(audio_clips)
                            aclip = aclip_looped.subclip(0, video_duration)
                        else:
                            # Audio rỗng, dùng silent
                            raise ValueError("Empty audio file")
                    elif audio_duration > video_duration:
                        # Audio dài hơn: cắt cho vừa
                        aclip = aclip.subclip(0, video_duration)
                    
                    # Gắn audio - bỏ qua kiểm tra to_soundarray vì nó có thể gây lỗi và không cần thiết
                    # Nếu file tồn tại, có size hợp lệ và duration > 0, thì có thể sử dụng
                    slide_clip = slide_clip.set_audio(aclip)
                    print(f"[AUDIO] Audio attached for slide {idx} (audio={audio_duration:.2f}s, video={video_duration:.2f}s, file_size={file_size} bytes)")
                except Exception as e:
                    print(f"[AUDIO] Failed to load audio: {e}, trying gTTS")
                    gtts_success = False
                    try:
                        from gtts import gTTS
                        if current_chunk_text:
                            gtts_path = os.path.join(tmpdir, f"gtts_{idx:03d}.mp3")
                            tts = gTTS(text=current_chunk_text[:500], lang='en', slow=False)
                            tts.save(gtts_path)
                            gtts_clip = AudioFileClip(gtts_path)
                            video_duration = slide_clip.duration
                            gtts_duration = gtts_clip.duration
                            
                            # Đảm bảo VIDEO và gTTS audio duration khớp nhau
                            if video_duration > gtts_duration:
                                # Video dài hơn: cắt video hoặc loop audio
                                slide_clip = slide_clip.subclip(0, gtts_duration)
                                video_duration = gtts_duration
                                aclip = gtts_clip
                                print(f"[DURATION FIX] Video cut to match gTTS: {video_duration:.2f}s")
                            elif gtts_duration > video_duration:
                                # Audio dài hơn: cắt audio
                                aclip = gtts_clip.subclip(0, video_duration)
                                print(f"[DURATION FIX] gTTS cut to match video: {video_duration:.2f}s")
                            else:
                                # Khớp nhau
                                aclip = gtts_clip
                            
                            slide_clip = slide_clip.set_audio(aclip)
                            print(f"[AUDIO] gTTS audio attached for slide {idx} (gTTS={gtts_duration:.2f}s, video={video_duration:.2f}s)")
                            gtts_success = True
                    except Exception as e2:
                        print(f"[AUDIO] gTTS also failed: {e2}")
                    if not gtts_success:
                        slide_clip = slide_clip.set_audio(_make_silent_audio(slide_clip.duration))
            else:
                gtts_success = False
                try:
                    from gtts import gTTS
                    if current_chunk_text:
                        gtts_path = os.path.join(tmpdir, f"gtts_{idx:03d}.mp3")
                        tts = gTTS(text=current_chunk_text[:500], lang='en', slow=False)
                        tts.save(gtts_path)
                        gtts_clip = AudioFileClip(gtts_path)
                        video_duration = slide_clip.duration
                        gtts_duration = gtts_clip.duration
                        
                        # Đảm bảo VIDEO và gTTS audio duration khớp nhau
                        if video_duration > gtts_duration:
                            # Video dài hơn: cắt video cho vừa với audio
                            slide_clip = slide_clip.subclip(0, gtts_duration)
                            video_duration = gtts_duration
                            aclip = gtts_clip
                            print(f"[DURATION FIX] Video cut to match gTTS: {video_duration:.2f}s")
                        elif gtts_duration > video_duration:
                            # Audio dài hơn: cắt audio cho vừa với video
                            aclip = gtts_clip.subclip(0, video_duration)
                            print(f"[DURATION FIX] gTTS cut to match video: {video_duration:.2f}s")
                        else:
                            # Khớp nhau
                            aclip = gtts_clip
                        
                        slide_clip = slide_clip.set_audio(aclip)
                        print(f"[AUDIO] gTTS audio created for slide {idx} (gTTS={gtts_duration:.2f}s, video={video_duration:.2f}s)")
                        gtts_success = True
                except Exception as e:
                    print(f"[AUDIO] gTTS failed: {e}")
                if not gtts_success:
                    slide_clip = slide_clip.set_audio(_make_silent_audio(slide_clip.duration))

            # Thêm fade in/out cho transitions mượt mà giữa các slide
            transition_duration = 0.8  # 0.8 giây fade để mượt hơn
            try:
                if cur_time > 0:  # Không fade in cho clip đầu tiên
                    slide_clip = fadein(slide_clip, transition_duration)
                slide_clip = fadeout(slide_clip, transition_duration)
                print(f"[TRANSITION] Added fade in/out for slide {idx}")
            except Exception as e:
                # Nếu fade không được hỗ trợ, bỏ qua
                print(f"[TRANSITION] Fade not supported: {e}, continuing without fade")

            clips.append(slide_clip)
            timeline.append((cur_time, cur_time + slide_clip.duration, current_chunk_text or ""))
            cur_time += slide_clip.duration

        if not clips:
            raise ValueError("Không tạo được clip nào.")

        final = concatenate_videoclips(clips, method="compose").set_fps(ENCODE_FPS)
        print(f"[ENCODE] Exporting video -> {out_video}")
        
        # Đảm bảo thư mục output tồn tại
        _safe_makedirs(os.path.dirname(out_video))
        
        # Đảm bảo đường dẫn không có ký tự đặc biệt gây lỗi
        out_video_clean = os.path.normpath(out_video)
        
        try:
            # Đơn giản hóa tối đa: để MoviePy tự quản lý tất cả temp files
            # Không chỉ định temp_audiofile để tránh lỗi với đường dẫn
            final.write_videofile(
                out_video_clean,
            codec="libx264",
            audio_codec="aac",
            preset="ultrafast",
                threads=4,
                bitrate="700k",
                fps=ENCODE_FPS,
            ffmpeg_params=["-movflags", "+faststart"],
                remove_temp=True,
            logger=None
        )
            print(f"[ENCODE] Video exported successfully")

        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            print(f"[ENCODE ERROR] {error_details}")
            raise RuntimeError(f"FFMPEG lỗi: {str(e)}")

        write_vtt(timeline, out_vtt)

        # Dọn dẹp
        for c in clips:
            _close_clip(c)
        _close_clip(final)
        gc.collect()

        return {
            "video_path": out_video,
            "caption_path": out_vtt,
            "script_text": "\n\n".join(chunks)
        }