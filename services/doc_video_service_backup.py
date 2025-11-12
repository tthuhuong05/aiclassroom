# services/doc_video_service.py
import os, uuid, textwrap, math, tempfile, contextlib, re
from typing import List, Tuple, Dict, Optional, Any
from PIL import Image, ImageDraw, ImageFont
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
import moviepy.video.fx.all as vfx
import moviepy.audio.fx.all as afx
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from services.ai_lecture_generator import ai_lecture_generator, LectureStructure, LectureSlide
from services.human_voice_service import human_voice_service, is_human_voice_available, synthesize_human_voice 
from services.image_service import image_service, search_images_for_content
from services.image_search_service import search_image, search_images, download_image
from services.luna_ai_service import luna_ai_service, is_luna_ai_available, analyze_content_for_luna_images
from services.alternative_image_service import alternative_image_service, is_alternative_image_available, analyze_content_for_alternative_images
from services.working_image_service import working_image_service, is_working_image_available, analyze_content_for_working_images
from services.kie_ai_service import kie_ai_service, is_kie_ai_available, analyze_content_for_kie_images
from PIL import ImageFilter

def _safe_makedirs(p):
    os.makedirs(p, exist_ok=True)

def _read_pdf(path: str) -> str:
    txt = []
    with open(path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            t = (page.extract_text() or "").strip()
            if t: txt.append(t)
    return "\n\n".join(txt)

def _read_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join((p.text or "").strip() for p in doc.paragraphs if (p.text or "").strip())

def _read_pptx(path: str) -> str:
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        buf = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                t = shape.text.strip()
                if t: buf.append(t)
        if buf: chunks.append("\n".join(buf))
    return "\n\n".join(chunks)

def extract_text_from_file(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf": return _read_pdf(path)
    if ext in (".docx",): return _read_docx(path)
    if ext in (".pptx", ".ppt"): return _read_pptx(path)
    if ext == ".txt": 
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    raise ValueError(f"Định dạng chưa hỗ trợ: {ext}")

def clean_text(s: str) -> str:
    s = re.sub(r'[ \t]+\n', '\n', s)
    s = re.sub(r'\n{3,}', '\n\n', s)
    return s.strip()

def chunk_text(s: str, max_chars=500) -> List[str]:
    # Ưu tiên ngắt theo câu; nếu dài thì wrap tiếp
    sentences = re.split(r'(?<=[\.\?\!])\s+', s)
    chunks, cur = [], ""
    for sent in sentences:
        if len(cur) + len(sent) + 1 <= max_chars:
            cur = (cur + " " + sent).strip()
        else:
            if cur: chunks.append(cur)
            if len(sent) <= max_chars:
                cur = sent
            else:
                # oversize sentence: hard wrap
                wraps = textwrap.wrap(sent, width=90)
                for w in wraps:
                    if len(w) > max_chars:
                        for sub in textwrap.wrap(w, width=70):
                            chunks.append(sub)
                        cur = ""
                    else:
                        chunks.append(w)
                cur = ""
    if cur: chunks.append(cur)
    # loại bỏ mẩu quá ngắn
    return [c.strip() for c in chunks if len(c.strip()) >= 20]

def estimate_duration_sec(text: str, wpm: int = 140) -> float:
    words = max(1, len(text.split()))
    return round(words / max(80, wpm) * 60.0, 2)

def synth_tts(text: str, out_wav: str, content_type: str = "general") -> float:
    """
    Trả về thời lượng giây với giọng đọc tự nhiên. Ưu tiên giọng người thật:
    1) Human voice (ElevenLabs) - Ưu tiên tuyệt đối giọng người thật
    2) pyttsx3 offline với cải tiến chất lượng
    3) gTTS online với cải tiến (chỉ khi bắt buộc)
    4) im lặng theo estimate
    """
    # Cải thiện text để giọng đọc tự nhiên hơn
    improved_text = _improve_text_for_tts(text, content_type)
    
    # 1) Human voice (ElevenLabs) - Ưu tiên tuyệt đối
    if is_human_voice_available():
        try:
            print("🎤 Using ElevenLabs human voice...")
            result = human_voice_service.synthesize_speech(improved_text, output_path=out_wav)
            if result.get("success"):
                print(f"✅ Human voice synthesis successful")
                return result.get("duration", estimate_duration_sec(improved_text))
            else:
                print(f"❌ Human voice synthesis failed: {result.get('error', 'Unknown error')}")
        except Exception as e:
            print(f"❌ Human voice synthesis error: {e}")
    
    # 2) pyttsx3 với cải tiến chất lượng cao
    try:
        import pyttsx3, time, wave
        print("Using pyttsx3 with enhanced quality...")
        eng = pyttsx3.init()
        
        # Cài đặt giọng đọc tự nhiên nhất có thể
        voices = eng.getProperty('voices')
        if voices:
            # Tìm giọng nữ tự nhiên nhất
            best_voice = None
            for voice in voices:
                voice_name = voice.name.lower()
                if any(keyword in voice_name for keyword in ['female', 'nữ', 'woman', 'zira', 'hazel']):
                    best_voice = voice
                    break
            
            if best_voice:
                eng.setProperty('voice', best_voice.id)
                print(f"✅ Using voice: {best_voice.name}")
        
        # Điều chỉnh tốc độ và âm lượng tối ưu
        rate = _get_optimal_rate(130, content_type)  # Chậm hơn để tự nhiên
        eng.setProperty('rate', rate)
        eng.setProperty('volume', 0.95)  # Âm lượng cao hơn
        
        eng.save_to_file(improved_text, out_wav)
        eng.runAndWait()
        
        # Đo duration thực tế
        if os.path.exists(out_wav):
            with contextlib.closing(wave.open(out_wav, 'r')) as wf:
                frames = wf.getnframes()
                rate = wf.getframerate()
                duration = frames / float(rate)
                print(f"✅ pyttsx3 synthesis completed: {duration:.1f}s")
                return duration
    except Exception as e:
        print(f"❌ pyttsx3 synthesis failed: {e}")
    
    # 3) gTTS với cải tiến (chỉ khi bắt buộc)
    try:
        from gtts import gTTS
        print("🌐 Using gTTS as fallback...")
        mp3_path = out_wav[:-4] + ".mp3"
        
        # Sử dụng gTTS với tùy chọn tối ưu nhất
        tts = gTTS(
            text=improved_text, 
            lang='vi', 
            slow=False,
            tld='com.vn'  # Domain Việt Nam
        )
        tts.save(mp3_path)
        
        # Convert và đo duration
        import math, mutagen
        if os.path.exists(mp3_path):
            dur = mutagen.File(mp3_path).info.length
            print(f"✅ gTTS synthesis completed: {dur:.1f}s")
            return float(dur)
    except Exception as e:
        print(f"❌ gTTS synthesis failed: {e}")
    
    # 4) Fallback - im lặng
    print("⚠️ All TTS methods failed, using estimated duration")
    return estimate_duration_sec(improved_text)

def _improve_text_for_tts(text: str, content_type: str) -> str:
    """Cải thiện text để giọng đọc tự nhiên hơn"""
    import re
    
    # Loại bỏ ký tự đặc biệt
    text = re.sub(r'[^\w\s.,!?;:()\-]', '', text)
    
    # Thêm dấu câu để tạo nhịp điệu tự nhiên
    text = re.sub(r'(\w+)\s+(\w+)', r'\1 \2', text)
    
    # Thêm dấu chấm câu cho câu dài
    sentences = text.split('.')
    improved_sentences = []
    for sentence in sentences:
        sentence = sentence.strip()
        if len(sentence) > 100:
            # Chia câu dài thành câu ngắn hơn
            words = sentence.split()
            mid_point = len(words) // 2
            first_part = ' '.join(words[:mid_point])
            second_part = ' '.join(words[mid_point:])
            improved_sentences.append(first_part + '. ' + second_part)
        else:
            improved_sentences.append(sentence)
    
    text = '. '.join(improved_sentences)
    
    # Thêm dấu câu cuối câu
    if not text.endswith(('.', '!', '?')):
        text += '.'
    
    return text

def _get_optimal_rate(wpm: int, content_type: str) -> int:
    """Tính toán tốc độ đọc tối ưu theo loại nội dung"""
    base_rate = wpm
    
    # Điều chỉnh tốc độ theo loại nội dung
    if content_type == "technical":
        return int(base_rate * 0.9)  # Chậm hơn cho nội dung kỹ thuật
    elif content_type == "educational":
        return int(base_rate * 0.95)  # Hơi chậm cho giáo dục
    elif content_type == "presentation":
        return int(base_rate * 1.05)  # Nhanh hơn cho thuyết trình
    else:
        return base_rate

def render_slide_image(text: str, w: int = 1280, h: int = 720, content_type: str = "general",
                       image_path: str = None) -> Image.Image:
    """Nếu image_path có, dùng ảnh làm nền mờ; nếu không vẫn sinh nền đẹp như hiện tại."""
    base = _create_dynamic_background(w, h, content_type)
    if image_path and os.path.exists(image_path):
        try:
            photo = Image.open(image_path).convert("RGB")
            # fill & crop theo khung
            pw, ph = photo.size
            scale = max(w/pw, h/ph)
            photo = photo.resize((int(pw*scale), int(ph*scale)))
            # crop trung tâm
            x = (photo.width - w)//2; y = (photo.height - h)//2
            photo = photo.crop((x, y, x+w, y+h)).filter(ImageFilter.GaussianBlur(6))
            base = Image.blend(photo, base, alpha=0.3)  # overlay nhẹ
        except Exception:
            pass

    draw = ImageDraw.Draw(base)
    try:
        font = ImageFont.truetype("arial.ttf", 48)
    except:
        font = ImageFont.load_default()

    lines = _improve_text_layout(text, w, font)
    y = 120
    for line in lines[:12]:
        if not line.strip(): y += 20; continue
        draw.text((102, y+2), line, fill=(0,0,0,100), font=font)
        draw.text((100, y),   line, fill=(255,255,255) if image_path else (22,27,34), font=font)
        y += 50

    _add_decorative_elements(draw, w, h, content_type)
    try:
        f2 = ImageFont.truetype("arial.ttf", 28)
    except:
        f2 = ImageFont.load_default()
    draw.rectangle([(0, h-80), (w, h)], fill=(0,0,0,120) if image_path else (229,231,235))
    draw.text((40, h-60), "🎬 Bài giảng tự động từ tài liệu", fill=(255,255,255) if image_path else (20,20,20), font=f2)
    return base


def _create_dynamic_background(w: int, h: int, content_type: str) -> Image.Image:
    """Tạo background động theo loại nội dung"""
    if content_type == "technical":
        # Background cho nội dung kỹ thuật
        img = Image.new("RGB", (w, h), color=(240, 248, 255))
        draw = ImageDraw.Draw(img)
        # Thêm pattern kỹ thuật
        for i in range(0, w, 40):
            for j in range(0, h, 40):
                draw.rectangle([i, j, i+20, j+20], fill=(230, 240, 255))
    elif content_type == "educational":
        # Background cho giáo dục
        img = Image.new("RGB", (w, h), color=(248, 250, 252))
        draw = ImageDraw.Draw(img)
        # Thêm pattern giáo dục
        for i in range(0, w, 60):
            for j in range(0, h, 60):
                draw.ellipse([i, j, i+30, j+30], fill=(240, 245, 250))
    elif content_type == "presentation":
        # Background cho thuyết trình
        img = Image.new("RGB", (w, h), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        # Thêm gradient effect
        for y in range(h):
            color_value = int(255 - (y / h) * 20)
            draw.line([(0, y), (w, y)], fill=(color_value, color_value, 255))
    else:
        # Background mặc định
        img = Image.new("RGB", (w, h), color=(248, 250, 252))
    
    return img

def _improve_text_layout(text: str, w: int, font) -> List[str]:
    """Cải thiện layout text"""
    # Chia text thành các đoạn
    paragraphs = text.split('\n\n')
    lines = []
    
    for paragraph in paragraphs:
        if not paragraph.strip():
            continue
            
        # Wrap text cho từng đoạn
        paragraph_lines = textwrap.wrap(paragraph, width=50)
        lines.extend(paragraph_lines)
        lines.append("")  # Thêm khoảng trống giữa các đoạn
    
    return lines

def _add_decorative_elements(draw, w: int, h: int, content_type: str):
    """Thêm các element trang trí"""
    if content_type == "technical":
        # Thêm icon kỹ thuật
        draw.rectangle([20, 20, 40, 40], fill=(59, 130, 246), outline=(37, 99, 235))
        draw.rectangle([w-40, 20, w-20, 40], fill=(59, 130, 246), outline=(37, 99, 235))
    elif content_type == "educational":
        # Thêm icon giáo dục
        draw.ellipse([20, 20, 50, 50], fill=(34, 197, 94), outline=(22, 163, 74))
        draw.ellipse([w-50, 20, w-20, 50], fill=(34, 197, 94), outline=(22, 163, 74))
    elif content_type == "presentation":
        # Thêm icon thuyết trình
        draw.polygon([(20, 30), (30, 20), (40, 30), (30, 40)], fill=(168, 85, 247), outline=(147, 51, 234))
        draw.polygon([(w-40, 30), (w-30, 20), (w-20, 30), (w-30, 40)], fill=(168, 85, 247), outline=(147, 51, 234))

def write_vtt(cues: List[Tuple[float, float, str]], outpath: str):
    def fmt(t):
        ms = int(round((t - int(t)) * 1000))
        s = int(t); h = s // 3600; m = (s % 3600) // 60; s = (s % 60)
        return f"{h:02d}:{m:02d}:{s:02d}.{ms:03d}"
    lines = ["WEBVTT", ""]
    for st, et, txt in cues:
        lines.append(f"{fmt(st)} --> {fmt(et)}")
        lines.append(txt.strip())
        lines.append("")
    _safe_makedirs(os.path.dirname(outpath))
    with open(outpath, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

def make_video_from_file(src_path: str, out_dir: str, title: str = None,
                         width=1280, height=720, wpm=140) -> Dict[str, str]:
    _safe_makedirs(out_dir)
    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)
    chunks = chunk_text(raw, max_chars=550)
    if not chunks: raise ValueError("Không trích xuất được nội dung đủ dài từ tài liệu.")

    clips = []
    timeline = []
    cur_t = 0.0
    tmpdir = tempfile.mkdtemp(prefix="doc2vid_")
    audio_paths = []

    for i, text in enumerate(chunks, 1):
        # tạo ảnh nền
        im = render_slide_image(text, w=width, h=height, content_type="educational")
        frame_path = os.path.join(tmpdir, f"frame_{i:03d}.png")
        im.save(frame_path)

        # tạo audio
        wav_path = os.path.join(tmpdir, f"tts_{i:03d}.wav")
        dur = synth_tts(text, wav_path, "educational")
        audio_path = wav_path if os.path.exists(wav_path) else None

        # nếu dùng gTTS mp3:
        if not audio_path:
            mp3_alt = wav_path[:-4] + ".mp3"
            if os.path.exists(mp3_alt):
                audio_path = mp3_alt

        # nếu không có file audio vật lý -> tạo ImageClip duration 'dur'
        clip = ImageClip(frame_path).set_duration(max(2.0, float(dur)))
        if audio_path:
            aclip = AudioFileClip(audio_path)
            clip = clip.set_audio(aclip).set_duration(aclip.duration)
            audio_paths.append(audio_path)

        # ghi timeline + phụ đề
        start = cur_t
        end = start + clip.duration
        timeline.append((start, end, text))
        cur_t = end

        clips.append(clip)

    final = concatenate_videoclips(clips, method="compose")
    base = (title or os.path.splitext(os.path.basename(src_path))[0]).strip() or "bai_giang"
    vid_name = f"{uuid.uuid4()}_{re.sub(r'[^a-zA-Z0-9_-]+','_', base)}.mp4"
    vtt_name = vid_name.replace(".mp4", ".vtt")

    out_video = os.path.join(out_dir, vid_name)
    out_vtt = os.path.join(out_dir, vtt_name)

    final.write_videofile(
        out_video, fps=24, codec="libx264", audio_codec="aac",
        temp_audiofile=os.path.join(tmpdir, "temp-audio.m4a"),
        remove_temp=True
    )
    write_vtt(timeline, out_vtt)

    return {
        "video_path": out_video,
        "caption_path": out_vtt,
        "script_text": "\n\n".join(chunks)
    }



def synth_voice(text: str, out_path: str) -> float:
    """
    Ưu tiên tuyệt đối giọng người thật (ElevenLabs). 
    Chỉ fallback về TTS máy khi ElevenLabs hoàn toàn không khả dụng.
    Trả về thời lượng ước tính hoặc thực tế (giây).
    """
    # Cho phép bắt buộc giọng người thật qua biến môi trường
    require_human = os.getenv('REQUIRE_HUMAN_VOICE', '0') == '1'
    if require_human and not is_human_voice_available():
        raise RuntimeError(
            "REQUIRE_HUMAN_VOICE=1 nhưng thiếu ELEVENLABS_API_KEY. "
            "Theo yêu cầu, không được dùng giọng máy."
        )

    print(f"🎤 Attempting human voice synthesis for: {text[:50]}...")
    
    # Ưu tiên tuyệt đối ElevenLabs
    if is_human_voice_available():
        try:
            print("✅ ElevenLabs API key found, using human voice...")
            res = synthesize_human_voice(text, output_path=out_path)
            if res.get("success") and os.path.exists(res["file_path"]):
                print(f"✅ Human voice synthesis successful: {res['file_path']}")
                return float(res.get("duration") or estimate_duration_sec(text))
            else:
                print(f"❌ Human voice synthesis failed: {res.get('error', 'Unknown error')}")
                # Thử lại với text được cải thiện
                improved_text = _improve_text_for_tts(text, "educational")
                if improved_text != text:
                    print("🔄 Retrying with improved text...")
                    res = synthesize_human_voice(improved_text, output_path=out_path)
                    if res.get("success") and os.path.exists(res["file_path"]):
                        print(f"✅ Human voice synthesis successful with improved text")
                        return float(res.get("duration") or estimate_duration_sec(improved_text))
        except Exception as e:
            print(f"❌ Human voice synthesis error: {e}")
    else:
        print("⚠️ ElevenLabs API key not found")
    
    # Fallback về TTS sẵn có (chỉ khi ElevenLabs hoàn toàn không khả dụng)
    print("🔄 Falling back to TTS...")
    return synth_tts(text, out_path, "educational")



def _build_image_queries_from_slide(slide) -> List[str]:
    """
    Cải thiện thuật toán tìm kiếm hình ảnh phù hợp với nội dung slide.
    Sử dụng AI để phân tích nội dung và đề xuất từ khóa hình ảnh chính xác.
    """
    candidates = []
    
    # Bước 1: Sử dụng AI để phân tích nội dung và đề xuất từ khóa hình ảnh
    try:
        from ai_gemini import _ensure
        if _ensure():
            ai_keywords = _get_ai_image_keywords(slide)
            if ai_keywords:
                candidates.extend(ai_keywords)
    except Exception as e:
        print(f"AI image keyword extraction failed: {e}")
    
    # Bước 2: Fallback - phân tích từ khóa thông minh từ nội dung
    if not candidates:
        candidates = _extract_smart_keywords_from_slide(slide)
    
    # Bước 3: Làm sạch và tối ưu từ khóa
    candidates = _optimize_image_keywords(candidates)
    
    return candidates[:5]  # Giới hạn 5 từ khóa tốt nhất

def _get_ai_image_keywords(slide) -> List[str]:
    """Sử dụng AI để phân tích nội dung và đề xuất từ khóa hình ảnh chính xác"""
    try:
        import google.generativeai as genai
        
        slide_content = f"""
        Tiêu đề: {getattr(slide, 'title', '')}
        Nội dung: {getattr(slide, 'content', '')}
        Điểm quan trọng: {', '.join(getattr(slide, 'key_points', []))}
        """
        
        prompt = f"""
        Phân tích nội dung slide bài giảng dưới đây và đề xuất 3-5 từ khóa tìm kiếm hình ảnh phù hợp:

        NỘI DUNG SLIDE:
        {slide_content.strip()}

        YÊU CẦU:
        - Đề xuất từ khóa cụ thể, sinh động để tìm hình ảnh minh họa
        - Hình ảnh phải phù hợp với nội dung slide, không chỉ là từ khóa chung chung
        - Ưu tiên từ khóa thể hiện khái niệm, ví dụ thực tế, biểu đồ, case study
        - Từ khóa phải dễ tìm hình ảnh trên Unsplash/Pexels
        - Tránh từ khóa quá trừu tượng hoặc không có hình ảnh minh họa

        Trả về danh sách từ khóa, mỗi từ khóa trên một dòng:
        """
        
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        
        if response.text:
            keywords = [line.strip() for line in response.text.split('\n') if line.strip()]
            return keywords[:5]
            
    except Exception as e:
        print(f"AI keyword extraction error: {e}")
    
    return []

def _extract_smart_keywords_from_slide(slide) -> List[str]:
    """Phân tích từ khóa thông minh từ nội dung slide"""
    candidates = []
    
    # Lấy title làm từ khóa chính
    if getattr(slide, "title", None):
        candidates.append(slide.title)
    
    # Phân tích content để tìm khái niệm chính
    content = getattr(slide, 'content', '')
    if content:
        # Tìm các khái niệm kỹ thuật, thuật ngữ chuyên môn
        import re
        technical_terms = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', content)
        candidates.extend(technical_terms[:3])
        
        # Tìm các từ khóa dài và có ý nghĩa
        stop_words = {"và","là","của","các","những","một","trong","với","cho","khi","để","như","về","từ","đó",
                     "the","and","for","with","from","this","that","these","those","có","được","này","đó"}
        
        words = re.findall(r"[A-Za-zÀ-ỹ0-9\-]+", content)
        meaningful_words = [w for w in words if len(w) >= 4 and w.lower() not in stop_words]
        
        # Ưu tiên từ khóa dài và có ý nghĩa
        for word in sorted(meaningful_words, key=len, reverse=True)[:5]:
            if word not in candidates:
                candidates.append(word)
    
    # Thêm key points
    if getattr(slide, "key_points", None):
        kp = [k for k in slide.key_points if k]
        if kp:
            candidates.append(" ".join(kp[:2]))

    return candidates

def _optimize_image_keywords(keywords: List[str]) -> List[str]:
    """Tối ưu từ khóa để tìm hình ảnh tốt hơn"""
    optimized = []
    
    for keyword in keywords:
        if not keyword or len(keyword.strip()) < 3:
            continue
            
        keyword = keyword.strip()
        
        # Cải thiện từ khóa để tìm hình ảnh tốt hơn
        improved_keyword = _improve_keyword_for_image_search(keyword)
        
        if improved_keyword and improved_keyword.lower() not in [k.lower() for k in optimized]:
            optimized.append(improved_keyword)
    
    return optimized

def _improve_keyword_for_image_search(keyword: str) -> str:
    """Cải thiện từ khóa để tìm hình ảnh phù hợp hơn"""
    # Mapping từ khóa chuyên môn sang từ khóa có hình ảnh
    keyword_mapping = {
        # Công nghệ thông tin
        "machine learning": "artificial intelligence technology",
        "deep learning": "neural network technology",
        "data science": "data analysis charts",
        "programming": "coding computer screen",
        "algorithm": "computer algorithm visualization",
        "database": "database management system",
        "network": "computer network diagram",
        "security": "cybersecurity protection",
        
        # Kinh doanh
        "marketing": "digital marketing strategy",
        "finance": "financial analysis charts",
        "management": "business management team",
        "strategy": "business strategy planning",
        "leadership": "team leadership meeting",
        
        # Giáo dục
        "education": "online learning platform",
        "teaching": "teacher classroom",
        "learning": "student studying",
        "research": "scientific research lab",
        
        # Khoa học
        "science": "scientific laboratory",
        "chemistry": "chemistry lab equipment",
        "physics": "physics experiment",
        "biology": "biology laboratory",
        "mathematics": "mathematical equations",
        
        # Y tế
        "medicine": "medical healthcare",
        "health": "healthcare professionals",
        "treatment": "medical treatment",
        "diagnosis": "medical diagnosis",
        
        # Nghệ thuật
        "art": "artistic creativity",
        "design": "graphic design work",
        "music": "musical instruments",
        "literature": "books reading",
    }
    
    keyword_lower = keyword.lower()
    
    # Kiểm tra mapping trực tiếp
    for key, value in keyword_mapping.items():
        if key in keyword_lower:
            return value
    
    # Nếu không có mapping, cải thiện từ khóa
    if len(keyword.split()) == 1:
        # Từ đơn - thêm context
        if keyword_lower in ["data", "analysis", "study", "research"]:
            return f"{keyword} professional work"
        elif keyword_lower in ["system", "process", "method"]:
            return f"{keyword} workflow diagram"
        elif keyword_lower in ["theory", "concept", "principle"]:
            return f"{keyword} explanation illustration"
    
    return keyword




def make_ai_lecture_video(src_path: str, out_dir: str, title: str = None,
                          width=1280, height=720, wpm=140,
                          use_human_voice: bool = True, include_images: bool = True) -> Dict[str, str]:

    """
    Tạo video bài giảng AI-powered từ tài liệu
    Tương tự NotebookLM - AI đọc, hiểu và tạo bài giảng
    """
    _safe_makedirs(out_dir)
    
    # Bước 1: Trích xuất nội dung từ tài liệu
    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)
    if not raw: raise ValueError("Không trích xuất được nội dung từ tài liệu.")
    
    # Bước 2: AI phân tích và tạo cấu trúc bài giảng
    try:
        lecture_structure, lecture_script = ai_lecture_generator.create_lecture_from_document(raw)
    except Exception as e:
        # Fallback về phương pháp cũ nếu AI không hoạt động
        return make_video_from_file(src_path, out_dir, title, width, height, wpm)
    
    # Bước 3: Tạo video từ cấu trúc bài giảng AI
    clips = []
    timeline = []
    cur_t = 0.0
    tmpdir = tempfile.mkdtemp(prefix="ai_lecture_")
    
    # Slide giới thiệu
    intro_slide = f"""
    {lecture_structure.title}

    {lecture_structure.introduction}

    Mục tiêu học tập:
    {chr(10).join(f"• {obj}" for obj in lecture_structure.learning_objectives)}
     """
    
    # Tạo slide giới thiệu
    intro_clip = _create_lecture_slide(intro_slide, tmpdir, 0, width, height, wpm)
    if intro_clip:
        clips.append(intro_clip)
        timeline.append((cur_t, cur_t + intro_clip.duration, intro_slide))
        cur_t += intro_clip.duration
    
    # Tạo slides chính
    for i, slide in enumerate(lecture_structure.slides, 1):
        slide_content = f"""
{slide.title}

{slide.content}

Điểm quan trọng:
{chr(10).join(f"• {point}" for point in slide.key_points)}
"""
        img_path = None
        if include_images:
            queries = _build_image_queries_from_slide(slide)
            for q in queries:
                # tìm 2–3 ứng viên/nguồn
                cands = search_images(q, count=3, orientation="landscape")
                for c in cands:
                    got = download_image(c)
                    if got.get("success") and os.path.exists(got.get("file_path", "")):
                        img_path = got["file_path"]
                        break
                if img_path:
                    print(f"🖼️  Using image for slide '{getattr(slide,'title','')}': {q}")
                    break
            if not img_path:
                # phao cuối
                img_path = search_image(getattr(slide, 'title', 'education') or "education")

        slide_clip = _create_lecture_slide(
            slide_content, tmpdir, i, width, height, wpm, img_path=img_path, use_human_voice=use_human_voice
        )
        if slide_clip:
            clips.append(slide_clip)
            timeline.append((cur_t, cur_t + slide_clip.duration, slide_content))
            cur_t += slide_clip.duration

    # Slide kết luận
    conclusion_slide = f"""
Kết luận

{lecture_structure.conclusion}

Cảm ơn bạn đã theo dõi bài giảng!
"""
    
    conclusion_clip = _create_lecture_slide(conclusion_slide, tmpdir, len(lecture_structure.slides) + 1, width, height, wpm)
    if conclusion_clip:
        clips.append(conclusion_clip)
        timeline.append((cur_t, cur_t + conclusion_clip.duration, conclusion_slide))
        cur_t += conclusion_clip.duration
    
    # Ghép video
    if not clips:
        raise ValueError("Không tạo được video từ bài giảng AI.")
    
    final = concatenate_videoclips(clips, method="compose")
    base = (title or lecture_structure.title or os.path.splitext(os.path.basename(src_path))[0]).strip() or "ai_lecture"
    vid_name = f"{uuid.uuid4()}_{re.sub(r'[^a-zA-Z0-9_-]+','_', base)}.mp4"
    vtt_name = vid_name.replace(".mp4", ".vtt")
    
    out_video = os.path.join(out_dir, vid_name)
    out_vtt = os.path.join(out_dir, vtt_name)
    
    final.write_videofile(
        out_video, fps=24, codec="libx264", audio_codec="aac",
        temp_audiofile=os.path.join(tmpdir, "temp-audio.m4a"),
        remove_temp=True
    )
    write_vtt(timeline, out_vtt)
    
    return {
        "video_path": out_video,
        "caption_path": out_vtt,
        "script_text": lecture_script,
        "lecture_structure": {
            "title": lecture_structure.title,
            "learning_objectives": lecture_structure.learning_objectives,
            "total_duration": lecture_structure.total_duration,
            "slide_count": len(lecture_structure.slides)
        }
    }

def _create_lecture_slide(content: str, tmpdir: str, slide_num: int, width: int, height: int, wpm: int,
                          img_path: str = None, use_human_voice: bool = True):
    try:
        frame_path = os.path.join(tmpdir, f"slide_{slide_num:03d}.png")
        im = render_slide_image(content, w=width, h=height, content_type="educational", image_path=img_path)
        im.save(frame_path)

        # tạo audio (ưu tiên giọng người)
        wav_path = os.path.join(tmpdir, f"tts_{slide_num:03d}.mp3")
        if use_human_voice:
            print(f"🎤 Creating human voice for slide {slide_num}...")
            dur = synth_voice(content, wav_path)
        else:
            print(f"Using TTS for slide {slide_num}...")
            dur = synth_tts(content, wav_path, "educational")
        audio_path = wav_path if os.path.exists(wav_path) else None
        if not audio_path:  # nhánh gTTS có thể tạo .wav
            alt = wav_path[:-4] + ".wav"
            if os.path.exists(alt): audio_path = alt

        clip = ImageClip(frame_path).set_duration(max(3.0, float(dur)))
        if audio_path:
            aclip = AudioFileClip(audio_path)
            clip = clip.set_audio(aclip).set_duration(aclip.duration)
        # Hiệu ứng mềm để giống video mẫu
        clip = clip.fx(vfx.fadein, 0.25).fx(vfx.fadeout, 0.25)
        if clip.audio:
            clip = clip.fx(afx.audio_fadein, 0.25).fx(afx.audio_fadeout, 0.25)
        return clip
    except Exception as e:
        print(f"Error creating slide {slide_num}: {e}")
        return None


def make_gemini_lecture_video(src_path: str, out_dir: str, title: str = None,
                               width=1280, height=720) -> Dict[str, str]:
    """
    Tạo video bài giảng từ nội dung chính của file sử dụng AI Gemini.
    Ưu điểm:
    - AI Gemini phân tích và trích xuất NỘI DUNG CHÍNH với độ chính xác cao
    - Giọng người thật (ElevenLabs) tự nhiên, sinh động
    - Hình ảnh phù hợp và minh họa rõ ràng từ Pexels/Unsplash
    - Script đọc tự nhiên như người thật giảng bài
    
    Args:
        src_path: Đường dẫn file tài liệu (PDF, DOCX, PPTX)
        out_dir: Thư mục output
        title: Tiêu đề video (optional)
        width: Chiều rộng video (default: 1280)
        height: Chiều cao video (default: 720)
    
    Returns:
        Dict chứa video_path, caption_path, script_text, lecture_info
    """
    from ai_gemini import extract_main_content_from_document
    
    # Kiểm tra bắt buộc giọng người thật
    require_human = os.getenv('REQUIRE_HUMAN_VOICE', '0') == '1'
    if require_human and not is_human_voice_available():
        raise RuntimeError('ElevenLabs chưa được cấu hình nhưng REQUIRE_HUMAN_VOICE=1; từ chối chuyển sang TTS máy.')
    
    _safe_makedirs(out_dir)
    
    # Bước 1: Trích xuất nội dung từ tài liệu
    print("Extracting content from document...")
    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)
    if not raw:
        raise ValueError("Không trích xuất được nội dung từ tài liệu.")
    
    # Bước 2: AI Gemini phân tích nội dung chính
    print("AI Gemini analyzing main content...")
    lecture_data = extract_main_content_from_document(raw)
    if not lecture_data:
        print("⚠️ Gemini analysis failed, falling back to standard method...")
        return make_ai_lecture_video(src_path, out_dir, title, width, height, use_human_voice=True, include_images=True)
    
    print(f"✅ Extracted main topic: {lecture_data.get('main_topic', 'N/A')}")
    print(f"✅ Core concepts: {', '.join(lecture_data.get('core_concepts', []))}")
    
    # Bước 3: Tạo video từ cấu trúc Gemini
    clips = []
    timeline = []
    cur_t = 0.0
    tmpdir = tempfile.mkdtemp(prefix="gemini_lecture_")
    
    # Slide giới thiệu
    intro_data = lecture_data.get("introduction", {})
    intro_script = intro_data.get("script", f"Chào mừng đến với bài giảng {lecture_data.get('lecture_title', 'AI')}")
    intro_keywords = intro_data.get("image_keywords", [])
    
    print("\n🎬 Creating introduction slide...")
    intro_img = None
    if intro_keywords:
        print(f"🔍 Searching image for: {intro_keywords[0]}")
        intro_img = search_image(intro_keywords[0])
    
    intro_clip = _create_gemini_slide(
        intro_script, intro_data.get("key_points", []),
        tmpdir, 0, width, height, img_path=intro_img
    )
    if intro_clip:
        clips.append(intro_clip)
        timeline.append((cur_t, cur_t + intro_clip.duration, intro_script))
        cur_t += intro_clip.duration
    
    # Các slides chính
    slides_data = lecture_data.get("slides", [])
    for i, slide_data in enumerate(slides_data, 1):
        print(f"\n🎬 Creating slide {i}/{len(slides_data)}: {slide_data.get('title', 'N/A')}")
        
        # Tìm hình ảnh phù hợp
        img_path = None
        img_keywords = slide_data.get("image_keywords", [])
        if img_keywords:
            print(f"🔍 Searching image for: {img_keywords[0]}")
            img_path = search_image(img_keywords[0])
            if img_path:
                print(f"✅ Found image: {img_path}")
        
        # Script cho slide
        slide_script = slide_data.get("script", slide_data.get("main_content", ""))
        
        slide_clip = _create_gemini_slide(
            slide_script, slide_data.get("key_points", []),
            tmpdir, i, width, height, img_path=img_path
        )
        
        if slide_clip:
            clips.append(slide_clip)
            timeline.append((cur_t, cur_t + slide_clip.duration, slide_script))
            cur_t += slide_clip.duration
    
    # Slide kết luận
    conclusion_data = lecture_data.get("conclusion", {})
    conclusion_script = conclusion_data.get("script", "Cảm ơn bạn đã theo dõi bài giảng!")
    conclusion_keywords = conclusion_data.get("image_keywords", [])
    
    print("\n🎬 Creating conclusion slide...")
    conclusion_img = None
    if conclusion_keywords:
        print(f"🔍 Searching image for: {conclusion_keywords[0]}")
        conclusion_img = search_image(conclusion_keywords[0])
    
    conclusion_clip = _create_gemini_slide(
        conclusion_script, conclusion_data.get("key_takeaways", []),
        tmpdir, len(slides_data) + 1, width, height, img_path=conclusion_img
    )
    if conclusion_clip:
        clips.append(conclusion_clip)
        timeline.append((cur_t, cur_t + conclusion_clip.duration, conclusion_script))
        cur_t += conclusion_clip.duration
    
    # Ghép video
    if not clips:
        raise ValueError("Không tạo được video từ bài giảng Gemini.")
    
    print("\n🎥 Rendering final video...")
    final = concatenate_videoclips(clips, method="compose")
    
    video_title = title or lecture_data.get("lecture_title", os.path.splitext(os.path.basename(src_path))[0])
    vid_name = f"{uuid.uuid4()}_{re.sub(r'[^a-zA-Z0-9_-]+','_', video_title)}.mp4"
    vtt_name = vid_name.replace(".mp4", ".vtt")
    
    out_video = os.path.join(out_dir, vid_name)
    out_vtt = os.path.join(out_dir, vtt_name)
    
    final.write_videofile(
        out_video, fps=24, codec="libx264", audio_codec="aac",
        temp_audiofile=os.path.join(tmpdir, "temp-audio.m4a"),
        remove_temp=True
    )
    write_vtt(timeline, out_vtt)
    
    # Tạo full script
    full_script = "\n\n".join([
        f"# {lecture_data.get('lecture_title', 'Bài giảng AI')}",
        f"\n## Giới thiệu\n{intro_script}",
        *[f"\n## {slide.get('title', f'Slide {i}')}\n{slide.get('script', '')}" 
          for i, slide in enumerate(slides_data, 1)],
        f"\n## Kết luận\n{conclusion_script}"
    ])
    
    print(f"\n✅ Video created successfully: {out_video}")
    print(f"✅ Duration: {cur_t:.1f} seconds")
    
    return {
        "video_path": out_video,
        "caption_path": out_vtt,
        "script_text": full_script,
        "lecture_info": {
            "title": lecture_data.get("lecture_title", "Bài giảng AI"),
            "main_topic": lecture_data.get("main_topic", "N/A"),
            "core_concepts": lecture_data.get("core_concepts", []),
            "total_slides": len(slides_data) + 2,  # intro + slides + conclusion
            "total_duration_seconds": cur_t
        }
    }

def _create_gemini_slide(script: str, key_points: List[str], tmpdir: str, 
                        slide_num: int, width: int, height: int, img_path: str = None):
    """
    Tạo slide từ script Gemini với giọng người thật và hình ảnh.
    """
    try:
        # Tạo nội dung hiển thị trên slide
        display_content = script
        if key_points:
            display_content += "\n\nĐiểm quan trọng:\n" + "\n".join(f"• {point}" for point in key_points)
        
        frame_path = os.path.join(tmpdir, f"slide_{slide_num:03d}.png")
        im = render_slide_image(display_content, w=width, h=height, 
                               content_type="educational", image_path=img_path)
        im.save(frame_path)
        
        # Tạo audio với giọng người thật
        audio_path = os.path.join(tmpdir, f"audio_{slide_num:03d}.mp3")
        print(f"🎤 Synthesizing human voice for slide {slide_num}...")
        dur = synth_voice(script, audio_path)
        
        # Kiểm tra audio file
        if not os.path.exists(audio_path):
            # Fallback to wav
            alt_audio_path = audio_path.replace(".mp3", ".wav")
            if os.path.exists(alt_audio_path):
                audio_path = alt_audio_path
            else:
                audio_path = None
        
        clip = ImageClip(frame_path).set_duration(max(3.0, float(dur)))
        if audio_path:
            aclip = AudioFileClip(audio_path)
            clip = clip.set_audio(aclip).set_duration(aclip.duration)
            print(f"✅ Slide {slide_num} duration: {aclip.duration:.1f}s")
        
        return clip
        
    except Exception as e:
        print(f"❌ Error creating Gemini slide {slide_num}: {e}")
        return None


def make_luna_lecture_video(src_path: str, out_dir: str, title: str = None,
                           width=1280, height=720) -> Dict[str, str]:
    """
    Tạo video bài giảng sử dụng LunaAI để tạo hình ảnh chân thật.
    Ưu điểm:
    - LunaAI tạo hình ảnh chân thật, phù hợp với nội dung
    - AI Gemini phân tích nội dung chuyên sâu
    - Giọng người thật (ElevenLabs) tự nhiên
    - Hình ảnh được tối ưu cho từng slide cụ thể
    
    Args:
        src_path: Đường dẫn file tài liệu (PDF, DOCX, PPTX)
        out_dir: Thư mục output
        title: Tiêu đề video (optional)
        width: Chiều rộng video (default: 1280)
        height: Chiều cao video (default: 720)
    
    Returns:
        Dict chứa video_path, caption_path, script_text, lecture_info
    """
    from ai_gemini import extract_main_content_from_document
    
    # Kiểm tra LunaAI
    if not is_luna_ai_available():
        print("⚠️ LunaAI chưa được cấu hình, fallback về phương pháp khác...")
        return make_gemini_lecture_video(src_path, out_dir, title, width, height)
    
    # Kiểm tra bắt buộc giọng người thật
    require_human = os.getenv('REQUIRE_HUMAN_VOICE', '0') == '1'
    if require_human and not is_human_voice_available():
        raise RuntimeError('ElevenLabs chưa được cấu hình nhưng REQUIRE_HUMAN_VOICE=1; từ chối chuyển sang TTS máy.')
    
    _safe_makedirs(out_dir)
    
    # Bước 1: Trích xuất nội dung từ tài liệu
    print("Extracting content from document...")
    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)
    if not raw:
        raise ValueError("Không trích xuất được nội dung từ tài liệu.")
    
    # Bước 2: AI Gemini phân tích nội dung chính
    print("AI Gemini analyzing main content...")
    lecture_data = extract_main_content_from_document(raw)
    if not lecture_data:
        print("⚠️ Gemini analysis failed, falling back to standard method...")
        return make_gemini_lecture_video(src_path, out_dir, title, width, height)
    
    print(f"✅ Extracted main topic: {lecture_data.get('main_topic', 'N/A')}")
    print(f"✅ Core concepts: {', '.join(lecture_data.get('core_concepts', []))}")
    
    # Bước 3: Tạo video từ cấu trúc Gemini với LunaAI
    clips = []
    timeline = []
    cur_t = 0.0
    tmpdir = tempfile.mkdtemp(prefix="luna_lecture_")
    
    # Slide giới thiệu
    intro_data = lecture_data.get("introduction", {})
    intro_script = intro_data.get("script", f"Chào mừng đến với bài giảng {lecture_data.get('lecture_title', 'AI')}")
    intro_keywords = intro_data.get("image_keywords", [])
    
    print("\n🎬 Creating introduction slide with LunaAI...")
    intro_img = None
    if intro_keywords:
        print(f"🎨 LunaAI generating image for: {intro_keywords[0]}")
        intro_result = luna_ai_service.generate_image(intro_keywords[0], style="educational")
        if intro_result and intro_result.get("success"):
            intro_img = intro_result["file_path"]
            print(f"✅ LunaAI image created: {intro_img}")
    
    intro_clip = _create_luna_slide(
        intro_script, intro_data.get("key_points", []),
        tmpdir, 0, width, height, img_path=intro_img
    )
    if intro_clip:
        clips.append(intro_clip)
        timeline.append((cur_t, cur_t + intro_clip.duration, intro_script))
        cur_t += intro_clip.duration
    
    # Các slides chính
    slides_data = lecture_data.get("slides", [])
    for i, slide_data in enumerate(slides_data, 1):
        print(f"\n🎬 Creating slide {i}/{len(slides_data)} with LunaAI: {slide_data.get('title', 'N/A')}")
        
        # Tạo hình ảnh với LunaAI
        img_path = None
        slide_content = f"{slide_data.get('title', '')} {slide_data.get('main_content', '')}"
        
        # Phân tích nội dung để tạo prompt hình ảnh
        print(f"🎨 LunaAI analyzing content for slide {i}...")
        image_prompts = analyze_content_for_luna_images(slide_content)
        
        # Thử tạo hình ảnh với các prompts
        for prompt in image_prompts[:3]:  # Thử tối đa 3 prompts
            print(f"🎨 LunaAI generating: {prompt[:50]}...")
            result = luna_ai_service.generate_image(prompt, style="educational")
            
            if result and result.get("success"):
                img_path = result["file_path"]
                print(f"✅ LunaAI image created: {img_path}")
                break
            else:
                print(f"❌ LunaAI failed for prompt: {prompt[:30]}...")
        
        # Script cho slide
        slide_script = slide_data.get("script", slide_data.get("main_content", ""))
        
        slide_clip = _create_luna_slide(
            slide_script, slide_data.get("key_points", []),
            tmpdir, i, width, height, img_path=img_path
        )
        
        if slide_clip:
            clips.append(slide_clip)
            timeline.append((cur_t, cur_t + slide_clip.duration, slide_script))
            cur_t += slide_clip.duration
    
    # Slide kết luận
    conclusion_data = lecture_data.get("conclusion", {})
    conclusion_script = conclusion_data.get("script", "Cảm ơn bạn đã theo dõi bài giảng!")
    conclusion_keywords = conclusion_data.get("image_keywords", [])
    
    print("\n🎬 Creating conclusion slide with LunaAI...")
    conclusion_img = None
    if conclusion_keywords:
        print(f"🎨 LunaAI generating image for: {conclusion_keywords[0]}")
        conclusion_result = luna_ai_service.generate_image(conclusion_keywords[0], style="educational")
        if conclusion_result and conclusion_result.get("success"):
            conclusion_img = conclusion_result["file_path"]
            print(f"✅ LunaAI image created: {conclusion_img}")
    
    conclusion_clip = _create_luna_slide(
        conclusion_script, conclusion_data.get("key_takeaways", []),
        tmpdir, len(slides_data) + 1, width, height, img_path=conclusion_img
    )
    if conclusion_clip:
        clips.append(conclusion_clip)
        timeline.append((cur_t, cur_t + conclusion_clip.duration, conclusion_script))
        cur_t += conclusion_clip.duration
    
    # Ghép video
    if not clips:
        raise ValueError("Không tạo được video từ bài giảng LunaAI.")
    
    print("\n🎥 Rendering final video...")
    final = concatenate_videoclips(clips, method="compose")
    
    video_title = title or lecture_data.get("lecture_title", os.path.splitext(os.path.basename(src_path))[0])
    vid_name = f"{uuid.uuid4()}_{re.sub(r'[^a-zA-Z0-9_-]+','_', video_title)}.mp4"
    vtt_name = vid_name.replace(".mp4", ".vtt")
    
    out_video = os.path.join(out_dir, vid_name)
    out_vtt = os.path.join(out_dir, vtt_name)
    
    final.write_videofile(
        out_video, fps=24, codec="libx264", audio_codec="aac",
        temp_audiofile=os.path.join(tmpdir, "temp-audio.m4a"),
        remove_temp=True
    )
    write_vtt(timeline, out_vtt)
    
    # Tạo full script
    full_script = "\n\n".join([
        f"# {lecture_data.get('lecture_title', 'Bài giảng LunaAI')}",
        f"\n## Giới thiệu\n{intro_script}",
        *[f"\n## {slide.get('title', f'Slide {i}')}\n{slide.get('script', '')}" 
          for i, slide in enumerate(slides_data, 1)],
        f"\n## Kết luận\n{conclusion_script}"
    ])
    
    print(f"\n✅ LunaAI video created successfully: {out_video}")
    print(f"✅ Duration: {cur_t:.1f} seconds")
    print(f"✅ Total slides: {len(slides_data) + 2}")
    
    return {
        "video_path": out_video,
        "caption_path": out_vtt,
        "script_text": full_script,
        "lecture_info": {
            "title": lecture_data.get("lecture_title", "Bài giảng LunaAI"),
            "main_topic": lecture_data.get("main_topic", "N/A"),
            "core_concepts": lecture_data.get("core_concepts", []),
            "total_slides": len(slides_data) + 2,  # intro + slides + conclusion
            "total_duration_seconds": cur_t,
            "image_generator": "LunaAI"
        }
    }

def _create_luna_slide(script: str, key_points: List[str], tmpdir: str, 
                      slide_num: int, width: int, height: int, img_path: str = None):
    """
    Tạo slide từ script với hình ảnh LunaAI và giọng người thật.
    """
    try:
        # Tạo nội dung hiển thị trên slide
        display_content = script
        if key_points:
            display_content += "\n\nĐiểm quan trọng:\n" + "\n".join(f"• {point}" for point in key_points)
        
        frame_path = os.path.join(tmpdir, f"luna_slide_{slide_num:03d}.png")
        im = render_slide_image(display_content, w=width, h=height, 
                               content_type="educational", image_path=img_path)
        im.save(frame_path)
        
        # Tạo audio với giọng người thật
        audio_path = os.path.join(tmpdir, f"luna_audio_{slide_num:03d}.mp3")
        print(f"🎤 Synthesizing human voice for LunaAI slide {slide_num}...")
        dur = synth_voice(script, audio_path)
        
        # Kiểm tra audio file
        if not os.path.exists(audio_path):
            # Fallback to wav
            alt_audio_path = audio_path.replace(".mp3", ".wav")
            if os.path.exists(alt_audio_path):
                audio_path = alt_audio_path
            else:
                audio_path = None
        
        clip = ImageClip(frame_path).set_duration(max(3.0, float(dur)))
        if audio_path:
            aclip = AudioFileClip(audio_path)
            clip = clip.set_audio(aclip).set_duration(aclip.duration)
            print(f"✅ LunaAI slide {slide_num} duration: {aclip.duration:.1f}s")
        
        return clip
        
    except Exception as e:
        print(f"Error creating LunaAI slide {slide_num}: {e}")
        return None


def make_smart_lecture_video(src_path: str, out_dir: str, title: str = None,
                            width=1280, height=720) -> Dict[str, str]:
    """
    Tạo video bài giảng thông minh với hình ảnh tốt nhất có thể.
    Ưu điểm:
    - Thử LunaAI trước (nếu có)
    - Fallback về các dịch vụ thay thế (Pexels, Unsplash, OpenAI DALL-E)
    - AI Gemini phân tích nội dung chuyên sâu
    - Giọng người thật (ElevenLabs) tự nhiên
    - Hình ảnh được tối ưu cho từng slide cụ thể
    
    Args:
        src_path: Đường dẫn file tài liệu (PDF, DOCX, PPTX)
        out_dir: Thư mục output
        title: Tiêu đề video (optional)
        width: Chiều rộng video (default: 1280)
        height: Chiều cao video (default: 720)
    
    Returns:
        Dict chứa video_path, caption_path, script_text, lecture_info
    """
    from ai_gemini import extract_main_content_from_document
    
    # Kiểm tra các dịch vụ hình ảnh
    luna_available = is_luna_ai_available()
    alternative_available = is_alternative_image_available()
    
    if not luna_available and not alternative_available:
        print("No image generation services available, falling back to standard method...")
        return make_gemini_lecture_video(src_path, out_dir, title, width, height)
    
    # Kiểm tra bắt buộc giọng người thật
    require_human = os.getenv('REQUIRE_HUMAN_VOICE', '0') == '1'
    if require_human and not is_human_voice_available():
        raise RuntimeError('ElevenLabs chưa được cấu hình nhưng REQUIRE_HUMAN_VOICE=1; từ chối chuyển sang TTS máy.')
    
    _safe_makedirs(out_dir)
    
    # Bước 1: Trích xuất nội dung từ tài liệu
    print("Extracting content from document...")
    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)
    if not raw:
        raise ValueError("Không trích xuất được nội dung từ tài liệu.")
    
    # Bước 2: AI Gemini phân tích nội dung chính
    print("AI Gemini analyzing main content...")
    lecture_data = extract_main_content_from_document(raw)
    if not lecture_data:
        print("Gemini analysis failed, falling back to standard method...")
        return make_gemini_lecture_video(src_path, out_dir, title, width, height)
    
    print(f"Extracted main topic: {lecture_data.get('main_topic', 'N/A')}")
    print(f"Core concepts: {', '.join(lecture_data.get('core_concepts', []))}")
    
    # Bước 3: Tạo video từ cấu trúc Gemini với hình ảnh thông minh
    clips = []
    timeline = []
    cur_t = 0.0
    tmpdir = tempfile.mkdtemp(prefix="smart_lecture_")
    
    # Slide giới thiệu
    intro_data = lecture_data.get("introduction", {})
    intro_script = intro_data.get("script", f"Chào mừng đến với bài giảng {lecture_data.get('lecture_title', 'AI')}")
    intro_keywords = intro_data.get("image_keywords", [])
    
    print("\nCreating introduction slide with smart image generation...")
    intro_img = _generate_smart_image(intro_keywords[0] if intro_keywords else "introduction", luna_available, alternative_available)
    
    intro_clip = _create_smart_slide(
        intro_script, intro_data.get("key_points", []),
        tmpdir, 0, width, height, img_path=intro_img
    )
    if intro_clip:
        clips.append(intro_clip)
        timeline.append((cur_t, cur_t + intro_clip.duration, intro_script))
        cur_t += intro_clip.duration
    
    # Các slides chính
    slides_data = lecture_data.get("slides", [])
    for i, slide_data in enumerate(slides_data, 1):
        print(f"\nCreating slide {i}/{len(slides_data)}: {slide_data.get('title', 'N/A')}")
        
        # Tạo hình ảnh thông minh
        img_path = None
        slide_content = f"{slide_data.get('title', '')} {slide_data.get('main_content', '')}"
        
        # Phân tích nội dung để tạo prompt hình ảnh
        print(f"Analyzing content for slide {i}...")
        image_prompts = _analyze_content_for_smart_images(slide_content, luna_available, alternative_available)
        
        # Thử tạo hình ảnh với các prompts
        for prompt in image_prompts[:3]:  # Thử tối đa 3 prompts
            print(f"Generating image: {prompt[:50]}...")
            img_path = _generate_smart_image(prompt, luna_available, alternative_available)
            
            if img_path:
                print(f"Image created: {img_path}")
                break
            else:
                print(f"Failed for prompt: {prompt[:30]}...")
        
        # Script cho slide
        slide_script = slide_data.get("script", slide_data.get("main_content", ""))
        
        slide_clip = _create_smart_slide(
            slide_script, slide_data.get("key_points", []),
            tmpdir, i, width, height, img_path=img_path
        )
        
        if slide_clip:
            clips.append(slide_clip)
            timeline.append((cur_t, cur_t + slide_clip.duration, slide_script))
            cur_t += slide_clip.duration
    
    # Slide kết luận
    conclusion_data = lecture_data.get("conclusion", {})
    conclusion_script = conclusion_data.get("script", "Cảm ơn bạn đã theo dõi bài giảng!")
    conclusion_keywords = conclusion_data.get("image_keywords", [])
    
    print("\nCreating conclusion slide with smart image generation...")
    conclusion_img = _generate_smart_image(conclusion_keywords[0] if conclusion_keywords else "conclusion", luna_available, alternative_available)
    
    conclusion_clip = _create_smart_slide(
        conclusion_script, conclusion_data.get("key_takeaways", []),
        tmpdir, len(slides_data) + 1, width, height, img_path=conclusion_img
    )
    if conclusion_clip:
        clips.append(conclusion_clip)
        timeline.append((cur_t, cur_t + conclusion_clip.duration, conclusion_script))
        cur_t += conclusion_clip.duration
    
    # Ghép video
    if not clips:
        raise ValueError("Không tạo được video từ bài giảng thông minh.")
    
    print("\nRendering final video...")
    final = concatenate_videoclips(clips, method="compose")
    
    video_title = title or lecture_data.get("lecture_title", os.path.splitext(os.path.basename(src_path))[0])
    vid_name = f"{uuid.uuid4()}_{re.sub(r'[^a-zA-Z0-9_-]+','_', video_title)}.mp4"
    vtt_name = vid_name.replace(".mp4", ".vtt")
    
    out_video = os.path.join(out_dir, vid_name)
    out_vtt = os.path.join(out_dir, vtt_name)
    
    final.write_videofile(
        out_video, fps=24, codec="libx264", audio_codec="aac",
        temp_audiofile=os.path.join(tmpdir, "temp-audio.m4a"),
        remove_temp=True
    )
    write_vtt(timeline, out_vtt)
    
    # Tạo full script
    full_script = "\n\n".join([
        f"# {lecture_data.get('lecture_title', 'Bài giảng thông minh')}",
        f"\n## Giới thiệu\n{intro_script}",
        *[f"\n## {slide.get('title', f'Slide {i}')}\n{slide.get('script', '')}" 
          for i, slide in enumerate(slides_data, 1)],
        f"\n## Kết luận\n{conclusion_script}"
    ])
    
    # Xác định dịch vụ hình ảnh được sử dụng
    image_service_used = "Unknown"
    if luna_available:
        image_service_used = "LunaAI"
    elif alternative_available:
        image_service_used = "Alternative Services"
    
    print(f"\nSmart video created successfully: {out_video}")
    print(f"Duration: {cur_t:.1f} seconds")
    print(f"Total slides: {len(slides_data) + 2}")
    print(f"Image service: {image_service_used}")
    
    return {
        "video_path": out_video,
        "caption_path": out_vtt,
        "script_text": full_script,
        "lecture_info": {
            "title": lecture_data.get("lecture_title", "Bài giảng thông minh"),
            "main_topic": lecture_data.get("main_topic", "N/A"),
            "core_concepts": lecture_data.get("core_concepts", []),
            "total_slides": len(slides_data) + 2,  # intro + slides + conclusion
            "total_duration_seconds": cur_t,
            "image_generator": image_service_used
        }
    }

def _generate_smart_image(prompt: str, luna_available: bool, alternative_available: bool) -> Optional[str]:
    """Tạo hình ảnh thông minh với fallback"""
    if luna_available:
        try:
            result = luna_ai_service.generate_image(prompt, style="educational")
            if result and result.get("success"):
                return result["file_path"]
        except Exception as e:
            print(f"LunaAI failed: {e}")
    
    if alternative_available:
        try:
            result = alternative_image_service.generate_image(prompt, style="educational")
            if result and result.get("success"):
                return result["file_path"]
        except Exception as e:
            print(f"Alternative service failed: {e}")
    
    return None

def _analyze_content_for_smart_images(content: str, luna_available: bool, alternative_available: bool) -> List[str]:
    """Phân tích nội dung để tạo prompts hình ảnh thông minh"""
    if luna_available:
        try:
            return analyze_content_for_luna_images(content)
        except Exception as e:
            print(f"LunaAI analysis failed: {e}")
    
    if alternative_available:
        try:
            return analyze_content_for_alternative_images(content)
        except Exception as e:
            print(f"Alternative analysis failed: {e}")
    
    # Fallback đơn giản
    import re
    words = re.findall(r'\b[A-Za-zÀ-ỹ0-9\-]+\b', content)
    important_words = [w for w in words if len(w) >= 4][:5]
    return [f"{word.lower()} professional" for word in important_words[:3]]

def _create_smart_slide(script: str, key_points: List[str], tmpdir: str, 
                       slide_num: int, width: int, height: int, img_path: str = None):
    """
    Tạo slide từ script với hình ảnh thông minh và giọng người thật.
    """
    try:
        # Tạo nội dung hiển thị trên slide
        display_content = script
        if key_points:
            display_content += "\n\nĐiểm quan trọng:\n" + "\n".join(f"• {point}" for point in key_points)
        
        frame_path = os.path.join(tmpdir, f"smart_slide_{slide_num:03d}.png")
        im = render_slide_image(display_content, w=width, h=height, 
                               content_type="educational", image_path=img_path)
        im.save(frame_path)
        
        # Tạo audio với giọng người thật
        audio_path = os.path.join(tmpdir, f"smart_audio_{slide_num:03d}.mp3")
        print(f"Synthesizing human voice for smart slide {slide_num}...")
        dur = synth_voice(script, audio_path)
        
        # Kiểm tra audio file
        if not os.path.exists(audio_path):
            # Fallback to wav
            alt_audio_path = audio_path.replace(".mp3", ".wav")
            if os.path.exists(alt_audio_path):
                audio_path = alt_audio_path
            else:
                audio_path = None
        
        clip = ImageClip(frame_path).set_duration(max(3.0, float(dur)))
        if audio_path:
            aclip = AudioFileClip(audio_path)
            clip = clip.set_audio(aclip).set_duration(aclip.duration)
            print(f"Smart slide {slide_num} duration: {aclip.duration:.1f}s")
        
        return clip
        
    except Exception as e:
        print(f"Error creating smart slide {slide_num}: {e}")
        return None


def make_working_lecture_video(src_path: str, out_dir: str, title: str = None,
                              width=1280, height=720) -> Dict[str, str]:
    """
    Tạo video bài giảng sử dụng các dịch vụ hình ảnh thực sự hoạt động.
    Ưu điểm:
    - Sử dụng Stability AI, OpenAI DALL-E, Pexels, Unsplash
    - AI Gemini phân tích nội dung chuyên sâu
    - Giọng người thật (ElevenLabs) tự nhiên
    - Hình ảnh chất lượng cao và phù hợp
    
    Args:
        src_path: Đường dẫn file tài liệu (PDF, DOCX, PPTX)
        out_dir: Thư mục output
        title: Tiêu đề video (optional)
        width: Chiều rộng video (default: 1280)
        height: Chiều cao video (default: 720)
    
    Returns:
        Dict chứa video_path, caption_path, script_text, lecture_info
    """
    from ai_gemini import extract_main_content_from_document
    
    # Kiểm tra các dịch vụ hình ảnh
    working_available = is_working_image_available()
    
    if not working_available:
        print("No working image services available, falling back to standard method...")
        return make_gemini_lecture_video(src_path, out_dir, title, width, height)
    
    # Kiểm tra bắt buộc giọng người thật
    require_human = os.getenv('REQUIRE_HUMAN_VOICE', '0') == '1'
    if require_human and not is_human_voice_available():
        raise RuntimeError('ElevenLabs chưa được cấu hình nhưng REQUIRE_HUMAN_VOICE=1; từ chối chuyển sang TTS máy.')
    
    _safe_makedirs(out_dir)
    
    # Bước 1: Trích xuất nội dung từ tài liệu
    print("Extracting content from document...")
    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)
    if not raw:
        raise ValueError("Không trích xuất được nội dung từ tài liệu.")
    
    # Bước 2: AI Gemini phân tích nội dung chính
    print("AI Gemini analyzing main content...")
    lecture_data = extract_main_content_from_document(raw)
    if not lecture_data:
        print("Gemini analysis failed, falling back to standard method...")
        return make_gemini_lecture_video(src_path, out_dir, title, width, height)
    
    print(f"Extracted main topic: {lecture_data.get('main_topic', 'N/A')}")
    print(f"Core concepts: {', '.join(lecture_data.get('core_concepts', []))}")
    
    # Bước 3: Tạo video từ cấu trúc Gemini với hình ảnh hoạt động
    clips = []
    timeline = []
    cur_t = 0.0
    tmpdir = tempfile.mkdtemp(prefix="working_lecture_")
    
    # Slide giới thiệu
    intro_data = lecture_data.get("introduction", {})
    intro_script = intro_data.get("script", f"Chào mừng đến với bài giảng {lecture_data.get('lecture_title', 'AI')}")
    intro_keywords = intro_data.get("image_keywords", [])
    
    print("\nCreating introduction slide with working image generation...")
    intro_img = _generate_working_image(intro_keywords[0] if intro_keywords else "introduction")
    
    intro_clip = _create_working_slide(
        intro_script, intro_data.get("key_points", []),
        tmpdir, 0, width, height, img_path=intro_img
    )
    if intro_clip:
        clips.append(intro_clip)
        timeline.append((cur_t, cur_t + intro_clip.duration, intro_script))
        cur_t += intro_clip.duration
    
    # Các slides chính
    slides_data = lecture_data.get("slides", [])
    for i, slide_data in enumerate(slides_data, 1):
        print(f"\nCreating slide {i}/{len(slides_data)}: {slide_data.get('title', 'N/A')}")
        
        # Tạo hình ảnh hoạt động
        img_path = None
        slide_content = f"{slide_data.get('title', '')} {slide_data.get('main_content', '')}"
        
        # Phân tích nội dung để tạo prompt hình ảnh
        print(f"Analyzing content for slide {i}...")
        image_prompts = analyze_content_for_working_images(slide_content)
        
        # Thử tạo hình ảnh với các prompts
        for prompt in image_prompts[:3]:  # Thử tối đa 3 prompts
            print(f"Generating image: {prompt[:50]}...")
            img_path = _generate_working_image(prompt)
            
            if img_path:
                print(f"Image created: {img_path}")
                break
            else:
                print(f"Failed for prompt: {prompt[:30]}...")
        
        # Script cho slide
        slide_script = slide_data.get("script", slide_data.get("main_content", ""))
        
        slide_clip = _create_working_slide(
            slide_script, slide_data.get("key_points", []),
            tmpdir, i, width, height, img_path=img_path
        )
        
        if slide_clip:
            clips.append(slide_clip)
            timeline.append((cur_t, cur_t + slide_clip.duration, slide_script))
            cur_t += slide_clip.duration
    
    # Slide kết luận
    conclusion_data = lecture_data.get("conclusion", {})
    conclusion_script = conclusion_data.get("script", "Cảm ơn bạn đã theo dõi bài giảng!")
    conclusion_keywords = conclusion_data.get("image_keywords", [])
    
    print("\nCreating conclusion slide with working image generation...")
    conclusion_img = _generate_working_image(conclusion_keywords[0] if conclusion_keywords else "conclusion")
    
    conclusion_clip = _create_working_slide(
        conclusion_script, conclusion_data.get("key_takeaways", []),
        tmpdir, len(slides_data) + 1, width, height, img_path=conclusion_img
    )
    if conclusion_clip:
        clips.append(conclusion_clip)
        timeline.append((cur_t, cur_t + conclusion_clip.duration, conclusion_script))
        cur_t += conclusion_clip.duration
    
    # Ghép video
    if not clips:
        raise ValueError("Không tạo được video từ bài giảng hoạt động.")
    
    print("\nRendering final video...")
    final = concatenate_videoclips(clips, method="compose")
    
    video_title = title or lecture_data.get("lecture_title", os.path.splitext(os.path.basename(src_path))[0])
    vid_name = f"{uuid.uuid4()}_{re.sub(r'[^a-zA-Z0-9_-]+','_', video_title)}.mp4"
    vtt_name = vid_name.replace(".mp4", ".vtt")
    
    out_video = os.path.join(out_dir, vid_name)
    out_vtt = os.path.join(out_dir, vtt_name)
    
    final.write_videofile(
        out_video, fps=24, codec="libx264", audio_codec="aac",
        temp_audiofile=os.path.join(tmpdir, "temp-audio.m4a"),
        remove_temp=True
    )
    write_vtt(timeline, out_vtt)
    
    # Tạo full script
    full_script = "\n\n".join([
        f"# {lecture_data.get('lecture_title', 'Bài giảng hoạt động')}",
        f"\n## Giới thiệu\n{intro_script}",
        *[f"\n## {slide.get('title', f'Slide {i}')}\n{slide.get('script', '')}" 
          for i, slide in enumerate(slides_data, 1)],
        f"\n## Kết luận\n{conclusion_script}"
    ])
    
    print(f"\nWorking video created successfully: {out_video}")
    print(f"Duration: {cur_t:.1f} seconds")
    print(f"Total slides: {len(slides_data) + 2}")
    print(f"Image service: Working Image Services")
    
    return {
        "video_path": out_video,
        "caption_path": out_vtt,
        "script_text": full_script,
        "lecture_info": {
            "title": lecture_data.get("lecture_title", "Bài giảng hoạt động"),
            "main_topic": lecture_data.get("main_topic", "N/A"),
            "core_concepts": lecture_data.get("core_concepts", []),
            "total_slides": len(slides_data) + 2,  # intro + slides + conclusion
            "total_duration_seconds": cur_t,
            "image_generator": "Working Image Services"
        }
    }

def _generate_working_image(prompt: str) -> Optional[str]:
    """Tạo hình ảnh sử dụng dịch vụ hoạt động"""
    try:
        result = working_image_service.generate_image(prompt, style="educational")
        if result and result.get("success"):
            return result["file_path"]
    except Exception as e:
        print(f"Working image service failed: {e}")
    
    return None

def _create_working_slide(script: str, key_points: List[str], tmpdir: str, 
                         slide_num: int, width: int, height: int, img_path: str = None):
    """
    Tạo slide từ script với hình ảnh hoạt động và giọng người thật.
    """
    try:
        # Tạo nội dung hiển thị trên slide
        display_content = script
        if key_points:
            display_content += "\n\nĐiểm quan trọng:\n" + "\n".join(f"• {point}" for point in key_points)
        
        frame_path = os.path.join(tmpdir, f"working_slide_{slide_num:03d}.png")
        im = render_slide_image(display_content, w=width, h=height, 
                               content_type="educational", image_path=img_path)
        im.save(frame_path)
        
        # Tạo audio với giọng người thật
        audio_path = os.path.join(tmpdir, f"working_audio_{slide_num:03d}.mp3")
        print(f"Synthesizing human voice for working slide {slide_num}...")
        dur = synth_voice(script, audio_path)
        
        # Kiểm tra audio file
        if not os.path.exists(audio_path):
            # Fallback to wav
            alt_audio_path = audio_path.replace(".mp3", ".wav")
            if os.path.exists(alt_audio_path):
                audio_path = alt_audio_path
            else:
                audio_path = None
        
        clip = ImageClip(frame_path).set_duration(max(3.0, float(dur)))
        if audio_path:
            aclip = AudioFileClip(audio_path)
            clip = clip.set_audio(aclip).set_duration(aclip.duration)
            print(f"Working slide {slide_num} duration: {aclip.duration:.1f}s")
        
        return clip
        
    except Exception as e:
        print(f"Error creating working slide {slide_num}: {e}")
        return None


def make_kie_lecture_video(src_path: str, out_dir: str, title: str = None,
                          width=1280, height=720) -> Dict[str, str]:
    """
    Tạo video bài giảng sử dụng KIE AI để tạo hình ảnh chân thật.
    Ưu điểm:
    - KIE AI tạo hình ảnh chân thật, phù hợp với nội dung
    - AI Gemini phân tích nội dung chuyên sâu
    - Giọng người thật (ElevenLabs) tự nhiên
    - Hình ảnh được tối ưu cho từng slide cụ thể
    
    Args:
        src_path: Đường dẫn file tài liệu (PDF, DOCX, PPTX)
        out_dir: Thư mục output
        title: Tiêu đề video (optional)
        width: Chiều rộng video (default: 1280)
        height: Chiều cao video (default: 720)
    
    Returns:
        Dict chứa video_path, caption_path, script_text, lecture_info
    """
    from ai_gemini import extract_main_content_from_document
    
    # Kiểm tra KIE AI
    if not is_kie_ai_available():
        print("KIE AI chưa được cấu hình, fallback về phương pháp khác...")
        return make_working_lecture_video(src_path, out_dir, title, width, height)
    
    # Kiểm tra bắt buộc giọng người thật
    require_human = os.getenv('REQUIRE_HUMAN_VOICE', '0') == '1'
    if require_human and not is_human_voice_available():
        raise RuntimeError('ElevenLabs chưa được cấu hình nhưng REQUIRE_HUMAN_VOICE=1; từ chối chuyển sang TTS máy.')
    
    _safe_makedirs(out_dir)
    
    # Bước 1: Trích xuất nội dung từ tài liệu
    print("Extracting content from document...")
    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)
    if not raw:
        raise ValueError("Không trích xuất được nội dung từ tài liệu.")
    
    # Bước 2: AI Gemini phân tích nội dung chính
    print("AI Gemini analyzing main content...")
    lecture_data = extract_main_content_from_document(raw)
    if not lecture_data:
        print("Gemini analysis failed, falling back to standard method...")
        return make_working_lecture_video(src_path, out_dir, title, width, height)
    
    print(f"Extracted main topic: {lecture_data.get('main_topic', 'N/A')}")
    print(f"Core concepts: {', '.join(lecture_data.get('core_concepts', []))}")
    
    # Bước 3: Tạo video từ cấu trúc Gemini với KIE AI
    clips = []
    timeline = []
    cur_t = 0.0
    tmpdir = tempfile.mkdtemp(prefix="kie_lecture_")
    
    # Slide giới thiệu
    intro_data = lecture_data.get("introduction", {})
    intro_script = intro_data.get("script", f"Chào mừng đến với bài giảng {lecture_data.get('lecture_title', 'AI')}")
    intro_keywords = intro_data.get("image_keywords", [])
    
    print("\nCreating introduction slide with KIE AI...")
    intro_img = None
    if intro_keywords:
        print(f"KIE AI generating image for: {intro_keywords[0]}")
        intro_result = kie_ai_service.generate_image(intro_keywords[0], style="educational")
        if intro_result and intro_result.get("success"):
            intro_img = intro_result["file_path"]
            print(f"KIE AI image created: {intro_img}")
    
    intro_clip = _create_kie_slide(
        intro_script, intro_data.get("key_points", []),
        tmpdir, 0, width, height, img_path=intro_img
    )
    if intro_clip:
        clips.append(intro_clip)
        timeline.append((cur_t, cur_t + intro_clip.duration, intro_script))
        cur_t += intro_clip.duration
    
    # Các slides chính
    slides_data = lecture_data.get("slides", [])
    for i, slide_data in enumerate(slides_data, 1):
        print(f"\nCreating slide {i}/{len(slides_data)} with KIE AI: {slide_data.get('title', 'N/A')}")
        
        # Tạo hình ảnh với KIE AI
        img_path = None
        slide_content = f"{slide_data.get('title', '')} {slide_data.get('main_content', '')}"
        
        # Phân tích nội dung để tạo prompt hình ảnh
        print(f"KIE AI analyzing content for slide {i}...")
        image_prompts = analyze_content_for_kie_images(slide_content)
        
        # Thử tạo hình ảnh với các prompts
        for prompt in image_prompts[:3]:  # Thử tối đa 3 prompts
            print(f"KIE AI generating: {prompt[:50]}...")
            result = kie_ai_service.generate_image(prompt, style="educational")
            
            if result and result.get("success"):
                img_path = result["file_path"]
                print(f"KIE AI image created: {img_path}")
                break
            else:
                print(f"KIE AI failed for prompt: {prompt[:30]}...")
        
        # Script cho slide
        slide_script = slide_data.get("script", slide_data.get("main_content", ""))
        
        slide_clip = _create_kie_slide(
            slide_script, slide_data.get("key_points", []),
            tmpdir, i, width, height, img_path=img_path
        )
        
        if slide_clip:
            clips.append(slide_clip)
            timeline.append((cur_t, cur_t + slide_clip.duration, slide_script))
            cur_t += slide_clip.duration
    
    # Slide kết luận
    conclusion_data = lecture_data.get("conclusion", {})
    conclusion_script = conclusion_data.get("script", "Cảm ơn bạn đã theo dõi bài giảng!")
    conclusion_keywords = conclusion_data.get("image_keywords", [])
    
    print("\nCreating conclusion slide with KIE AI...")
    conclusion_img = None
    if conclusion_keywords:
        print(f"KIE AI generating image for: {conclusion_keywords[0]}")
        conclusion_result = kie_ai_service.generate_image(conclusion_keywords[0], style="educational")
        if conclusion_result and conclusion_result.get("success"):
            conclusion_img = conclusion_result["file_path"]
            print(f"KIE AI image created: {conclusion_img}")
    
    conclusion_clip = _create_kie_slide(
        conclusion_script, conclusion_data.get("key_takeaways", []),
        tmpdir, len(slides_data) + 1, width, height, img_path=conclusion_img
    )
    if conclusion_clip:
        clips.append(conclusion_clip)
        timeline.append((cur_t, cur_t + conclusion_clip.duration, conclusion_script))
        cur_t += conclusion_clip.duration
    
    # Ghép video
    if not clips:
        raise ValueError("Không tạo được video từ bài giảng KIE AI.")
    
    print("\nRendering final video...")
    final = concatenate_videoclips(clips, method="compose")
    
    video_title = title or lecture_data.get("lecture_title", os.path.splitext(os.path.basename(src_path))[0])
    vid_name = f"{uuid.uuid4()}_{re.sub(r'[^a-zA-Z0-9_-]+','_', video_title)}.mp4"
    vtt_name = vid_name.replace(".mp4", ".vtt")
    
    out_video = os.path.join(out_dir, vid_name)
    out_vtt = os.path.join(out_dir, vtt_name)
    
    final.write_videofile(
        out_video, fps=24, codec="libx264", audio_codec="aac",
        temp_audiofile=os.path.join(tmpdir, "temp-audio.m4a"),
        remove_temp=True
    )
    write_vtt(timeline, out_vtt)
    
    # Tạo full script
    full_script = "\n\n".join([
        f"# {lecture_data.get('lecture_title', 'Bài giảng KIE AI')}",
        f"\n## Giới thiệu\n{intro_script}",
        *[f"\n## {slide.get('title', f'Slide {i}')}\n{slide.get('script', '')}" 
          for i, slide in enumerate(slides_data, 1)],
        f"\n## Kết luận\n{conclusion_script}"
    ])
    
    print(f"\nKIE AI video created successfully: {out_video}")
    print(f"Duration: {cur_t:.1f} seconds")
    print(f"Total slides: {len(slides_data) + 2}")
    
    return {
        "video_path": out_video,
        "caption_path": out_vtt,
        "script_text": full_script,
        "lecture_info": {
            "title": lecture_data.get("lecture_title", "Bài giảng KIE AI"),
            "main_topic": lecture_data.get("main_topic", "N/A"),
            "core_concepts": lecture_data.get("core_concepts", []),
            "total_slides": len(slides_data) + 2,  # intro + slides + conclusion
            "total_duration_seconds": cur_t,
            "image_generator": "KIE AI"
        }
    }

def _create_kie_slide(script: str, key_points: List[str], tmpdir: str, 
                     slide_num: int, width: int, height: int, img_path: str = None):
    """
    Tạo slide từ script với hình ảnh KIE AI và giọng người thật.
    """
    try:
        # Tạo nội dung hiển thị trên slide
        display_content = script
        if key_points:
            display_content += "\n\nĐiểm quan trọng:\n" + "\n".join(f"• {point}" for point in key_points)
        
        frame_path = os.path.join(tmpdir, f"kie_slide_{slide_num:03d}.png")
        im = render_slide_image(display_content, w=width, h=height, 
                               content_type="educational", image_path=img_path)
        im.save(frame_path)
        
        # Tạo audio với giọng người thật
        audio_path = os.path.join(tmpdir, f"kie_audio_{slide_num:03d}.mp3")
        print(f"Synthesizing human voice for KIE AI slide {slide_num}...")
        dur = synth_voice(script, audio_path)
        
        # Kiểm tra audio file
        if not os.path.exists(audio_path):
            # Fallback to wav
            alt_audio_path = audio_path.replace(".mp3", ".wav")
            if os.path.exists(alt_audio_path):
                audio_path = alt_audio_path
            else:
                audio_path = None
        
        clip = ImageClip(frame_path).set_duration(max(3.0, float(dur)))
        if audio_path:
            aclip = AudioFileClip(audio_path)
            clip = clip.set_audio(aclip).set_duration(aclip.duration)
            print(f"KIE AI slide {slide_num} duration: {aclip.duration:.1f}s")
        
        return clip
        
    except Exception as e:
        print(f"Error creating KIE AI slide {slide_num}: {e}")
        return None


def make_optimized_lecture_video(src_path: str, out_dir: str, title: str = None,
                                width: int = 1280, height: int = 720) -> Dict[str, str]:
    """
    Tao video bai giang toi uu su dung TAT CA API keys co san.
    Thu tu uu tien:
    1. Stability AI (tao hinh anh AI chat luong cao)
    2. Pexels (tim hinh anh phu hop)
    3. Unsplash (backup hinh anh)
    4. ElevenLabs (giong nguoi that)
    5. Gemini (phan tich noi dung chuyen sau)
    """
    from ai_gemini import extract_main_content_from_document
    
    # Kiem tra tat ca API keys
    apis_available = {
        "stability": bool(os.getenv("STABILITY_API_KEY")),
        "pexels": bool(os.getenv("PEXELS_API_KEY")), 
        "unsplash": bool(os.getenv("UNSPLASH_ACCESS_KEY")),
        "elevenlabs": is_human_voice_available(),
        "gemini": bool(os.getenv("GEMINI_API_KEY"))
    }
    
    print(f"APIs available: {[k for k, v in apis_available.items() if v]}")
    
    # Kiem tra bat buoc
    if not apis_available["elevenlabs"]:
        raise RuntimeError("ElevenLabs API key required for human voice")
    
    if not apis_available["gemini"]:
        raise RuntimeError("Gemini API key required for content analysis")
    
    if not any([apis_available["stability"], apis_available["pexels"], apis_available["unsplash"]]):
        raise RuntimeError("At least one image API required")
    
    _safe_makedirs(out_dir)
    
    # Buoc 1: Trich xuat noi dung
    print("Extracting content from document...")
    raw = extract_text_from_file(src_path)
    raw = clean_text(raw)
    if not raw:
        raise ValueError("Cannot extract content from document")
    
    # Buoc 2: AI phan tich noi dung
    print("AI analyzing content...")
    lecture_data = extract_main_content_from_document(raw)
    if not lecture_data:
        print("Gemini analysis failed, using fallback...")
        lecture_data = create_fallback_lecture_data(raw)
    
    print(f"Main topic: {lecture_data.get('main_topic', 'N/A')}")
    print(f"Core concepts: {', '.join(lecture_data.get('core_concepts', []))}")
    
    # Buoc 3: Tao video voi tat ca APIs
    clips = []
    timeline = []
    cur_t = 0.0
    tmpdir = tempfile.mkdtemp(prefix="optimized_lecture_")
    
    # Slide gioi thieu
    intro_data = lecture_data.get("introduction", {})
    intro_script = intro_data.get("script", f"Welcome to {lecture_data.get('lecture_title', 'AI Lecture')}")
    
    print("\nCreating introduction slide...")
    intro_img = generate_optimal_image(
        intro_data.get("image_keywords", ["introduction"]),
        apis_available
    )
    
    intro_clip = create_optimized_slide(
        intro_script, intro_data.get("key_points", []),
        tmpdir, 0, width, height, img_path=intro_img,
        apis_available
    )
    if intro_clip:
        clips.append(intro_clip)
        timeline.append((cur_t, cur_t + intro_clip.duration, intro_script))
        cur_t += intro_clip.duration
    
    # Cac slides chinh
    slides_data = lecture_data.get("slides", [])
    for i, slide_data in enumerate(slides_data, 1):
        print(f"\nCreating slide {i}/{len(slides_data)}: {slide_data.get('title', 'N/A')}")
        
        slide_content = f"{slide_data.get('title', '')} {slide_data.get('main_content', '')}"
        
        # Tao hinh anh toi uu
        print(f"Generating optimal image for slide {i}...")
        slide_img = generate_optimal_image(
            slide_data.get("image_keywords", []),
            apis_available
        )
        
        slide_script = slide_data.get("script", slide_data.get("main_content", ""))
        
        slide_clip = create_optimized_slide(
            slide_script, slide_data.get("key_points", []),
            tmpdir, i, width, height, img_path=slide_img,
            apis_available
        )
        
        if slide_clip:
            clips.append(slide_clip)
            timeline.append((cur_t, cur_t + slide_clip.duration, slide_script))
            cur_t += slide_clip.duration
    
    # Slide ket luan
    conclusion_data = lecture_data.get("conclusion", {})
    conclusion_script = conclusion_data.get("script", "Thank you for watching!")
    
    print("\nCreating conclusion slide...")
    conclusion_img = generate_optimal_image(
        conclusion_data.get("image_keywords", ["conclusion"]),
        apis_available
    )
    
    conclusion_clip = create_optimized_slide(
        conclusion_script, conclusion_data.get("key_takeaways", []),
        tmpdir, len(slides_data) + 1, width, height, img_path=conclusion_img,
        apis_available
    )
    if conclusion_clip:
        clips.append(conclusion_clip)
        timeline.append((cur_t, cur_t + conclusion_clip.duration, conclusion_script))
        cur_t += conclusion_clip.duration
    
    # Ghep video
    if not clips:
        raise ValueError("Cannot create video from lecture")
    
    print("\nRendering final video...")
    final = concatenate_videoclips(clips, method="compose")
    
    video_title = title or lecture_data.get("lecture_title", os.path.splitext(os.path.basename(src_path))[0])
    vid_name = f"{uuid.uuid4()}_{re.sub(r'[^a-zA-Z0-9_-]+','_', video_title)}.mp4"
    vtt_name = vid_name.replace(".mp4", ".vtt")
    
    out_video = os.path.join(out_dir, vid_name)
    out_vtt = os.path.join(out_dir, vtt_name)
    
    final.write_videofile(
        out_video, fps=24, codec="libx264", audio_codec="aac",
        temp_audiofile=os.path.join(tmpdir, "temp-audio.m4a"),
        remove_temp=True
    )
    write_vtt(timeline, out_vtt)
    
    # Tao full script
    full_script = "\n\n".join([
        f"# {lecture_data.get('lecture_title', 'Optimized AI Lecture')}",
        f"\n## Introduction\n{intro_script}",
        *[f"\n## {slide.get('title', f'Slide {i}')}\n{slide.get('script', '')}" 
          for i, slide in enumerate(slides_data, 1)],
        f"\n## Conclusion\n{conclusion_script}"
    ])
    
    print(f"\nOptimized video created successfully: {out_video}")
    print(f"Duration: {cur_t:.1f} seconds")
    print(f"Total slides: {len(slides_data) + 2}")
    print(f"APIs used: {[k for k, v in apis_available.items() if v]}")
    
    return {
        "video_path": out_video,
        "caption_path": out_vtt,
        "script_text": full_script,
        "lecture_info": {
            "title": lecture_data.get("lecture_title", "Optimized AI Lecture"),
            "main_topic": lecture_data.get("main_topic", "N/A"),
            "core_concepts": lecture_data.get("core_concepts", []),
            "total_slides": len(slides_data) + 2,
            "total_duration_seconds": cur_t,
            "apis_used": [k for k, v in apis_available.items() if v],
            "image_generator": "Optimized Multi-API"
        }
    }

def generate_optimal_image(keywords, apis_available):
    """Generate optimal image using best available API"""
    if not keywords:
        return None
    
    # Thu tu uu tien: Stability AI -> Pexels -> Unsplash
    if apis_available["stability"]:
        try:
            from services.working_image_service import generate_working_image
            result = generate_working_image(keywords[0], style="educational")
            if result and result.get("success"):
                print(f"   OK Stability AI: {result['file_path']}")
                return result["file_path"]
        except Exception as e:
            print(f"   LOI Stability AI: {e}")
    
    if apis_available["pexels"]:
        try:
            from services.working_image_service import generate_working_image
            result = generate_working_image(keywords[0], style="educational")
            if result and result.get("success"):
                print(f"   OK Pexels: {result['file_path']}")
                return result["file_path"]
        except Exception as e:
            print(f"   LOI Pexels: {e}")
    
    if apis_available["unsplash"]:
        try:
            from services.working_image_service import generate_working_image
            result = generate_working_image(keywords[0], style="educational")
            if result and result.get("success"):
                print(f"   OK Unsplash: {result['file_path']}")
                return result["file_path"]
        except Exception as e:
            print(f"   LOI Unsplash: {e}")
    
    print("   LOI: All image APIs failed")
    return None

def create_optimized_slide(script, key_points, tmpdir, slide_num, width, height, img_path, apis_available):
    """Create optimized slide with best available APIs"""
    try:
        # Tao noi dung hien thi
        display_content = script
        if key_points:
            display_content += "\n\nKey Points:\n" + "\n".join(f"• {point}" for point in key_points)
        
        frame_path = os.path.join(tmpdir, f"optimized_slide_{slide_num:03d}.png")
        im = render_slide_image(display_content, w=width, h=height, 
                               content_type="educational", image_path=img_path)
        im.save(frame_path)
        
        # Tao audio voi ElevenLabs
        audio_path = os.path.join(tmpdir, f"optimized_audio_{slide_num:03d}.mp3")
        print(f"   Synthesizing voice for slide {slide_num}...")
        dur = synth_voice(script, audio_path)
        
        # Kiem tra audio file
        if not os.path.exists(audio_path):
            alt_audio_path = audio_path.replace(".mp3", ".wav")
            if os.path.exists(alt_audio_path):
                audio_path = alt_audio_path
            else:
                audio_path = None
        
        clip = ImageClip(frame_path).set_duration(max(3.0, float(dur)))
        if audio_path:
            aclip = AudioFileClip(audio_path)
            clip = clip.set_audio(aclip).set_duration(aclip.duration)
            print(f"   Slide {slide_num} duration: {aclip.duration:.1f}s")
        
        return clip
        
    except Exception as e:
        print(f"   Error creating optimized slide {slide_num}: {e}")
        return None

def create_fallback_lecture_data(content):
    """Create fallback lecture data if Gemini fails"""
    return {
        "lecture_title": "AI Generated Lecture",
        "main_topic": "Educational Content",
        "core_concepts": ["Learning", "Education", "Knowledge"],
        "introduction": {
            "script": "Welcome to this educational presentation",
            "key_points": ["Introduction", "Overview"],
            "image_keywords": ["education", "learning"]
        },
        "slides": [
            {
                "title": "Main Content",
                "main_content": content[:200] + "..." if len(content) > 200 else content,
                "key_points": ["Key Point 1", "Key Point 2"],
                "script": content[:100] + "..." if len(content) > 100 else content,
                "image_keywords": ["education", "presentation"]
            }
        ],
        "conclusion": {
            "script": "Thank you for watching this presentation",
            "key_takeaways": ["Summary", "Conclusion"],
            "image_keywords": ["conclusion", "thank you"]
        }
    }
